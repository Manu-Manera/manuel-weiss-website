"""
Valkeen-gebrandete, 1:1-Nachbauten der drei Tempus/MS-Project-Screenshots.

Ausgabe:
  1) Valkeen_Integration-Optionen.pptx            – eigenständige, editierbare Datei
  2) Entscheidungsvorlage_Knauf_mit_Integrationsfolien.pptx
     – identisch mit der Original-Knauf-Vorlage, nur am Ende um die neuen
       Integrations-Folien erweitert (behält Master, Fonts, Footer und
       'Copyright Valkeen GmbH'-Branding).

Stil-Quelle: '/Users/manumanera/Downloads/Entscheidungsvorlage Knauf.pptx'
  - Navy:   #002856 (Valkeen-Primär)
  - Slate:  #335B74
  - Cyan:   #25B0CF
  - Green:  #86BC25
  - Orange: #F5891F
  - Red:    #C00000
  - Grey:   #B3B3B3 / #EFEFEF
  - Fonts:  'Calibri Light' (Titel), 'Calibri' (Fließtext)
"""

import copy
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Valkeen Palette
# ---------------------------------------------------------------------------
NAVY        = RGBColor(0x00, 0x28, 0x56)
NAVY_80     = RGBColor(0x33, 0x54, 0x78)  # 80 % Navy
SLATE       = RGBColor(0x33, 0x5B, 0x74)
CYAN        = RGBColor(0x25, 0xB0, 0xCF)
CYAN_LIGHT  = RGBColor(0xD4, 0xEF, 0xF6)
GREEN       = RGBColor(0x86, 0xBC, 0x25)
GREEN_LIGHT = RGBColor(0xE7, 0xF1, 0xCF)
ORANGE      = RGBColor(0xF5, 0x89, 0x1F)
ORANGE_LT   = RGBColor(0xFD, 0xE4, 0xCC)
RED         = RGBColor(0xC0, 0x00, 0x00)
GREY_40     = RGBColor(0xB3, 0xB3, 0xB3)
GREY_20     = RGBColor(0xD0, 0xD0, 0xCE)
GREY_10     = RGBColor(0xEF, 0xEF, 0xEF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
TEXT_DARK   = RGBColor(0x1C, 0x28, 0x35)
TEXT_MUTED  = RGBColor(0x59, 0x59, 0x59)

TITLE_FONT = "Calibri Light"
BODY_FONT  = "Calibri"

# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_fill(shape):
    shape.fill.background()


def _line(shape, color, width_pt=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def _no_line(shape):
    shape.line.fill.background()


def _kill_shadow(shape):
    spPr = shape._element.spPr
    # remove existing effectLst/effectRef
    for tag in ('a:effectLst', 'a:effectRef'):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    # append empty effectLst so master shadow can’t kick in
    effectLst = etree.SubElement(spPr, qn('a:effectLst'))


def _text(shape, text, *, size=14, bold=False, color=TEXT_DARK,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
          font=BODY_FONT, italic=False, line_spacing=None):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    tf.clear()
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color


def add_rect(slide, x, y, w, h, *, fill=WHITE, line=GREY_20, line_w=0.75,
             rounded=False, radius=0.08):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    if rounded:
        s.adjustments[0] = radius
    if fill is None:
        _no_fill(s)
    else:
        _fill(s, fill)
    if line is None:
        _no_line(s)
    else:
        _line(s, line, line_w)
    _kill_shadow(s)
    return s


def add_pill(slide, x, y, w, h, text, *, fill=ORANGE, text_color=WHITE,
             size=11, bold=True, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.adjustments[0] = 0.5  # full pill
    _fill(s, fill)
    if line is None:
        _no_line(s)
    else:
        _line(s, line, 0.75)
    _kill_shadow(s)
    _text(s, text, size=size, bold=bold, color=text_color, font=BODY_FONT)
    return s


def add_text(slide, x, y, w, h, text, *, size=12, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=BODY_FONT,
             italic=False, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _text(tb, text, size=size, bold=bold, color=color, align=align,
          anchor=anchor, font=font, italic=italic, line_spacing=line_spacing)
    return tb


def add_right_arrow_block(slide, x, y, w, h, *, fill, line=None,
                          head_ratio=0.10):
    """Breiter Pfeil als Block – MSO_SHAPE.RIGHT_ARROW."""
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               Inches(x), Inches(y),
                               Inches(w), Inches(h))
    # adjustment[0] = neck thickness (0..0.5), adjustment[1] = head width (0..0.5)
    try:
        s.adjustments[0] = 0.55  # dicker Schaft (gute Höhe für Pills)
        s.adjustments[1] = head_ratio
    except Exception:
        pass
    _fill(s, fill)
    if line is None:
        _no_line(s)
    else:
        _line(s, line, 0.75)
    _kill_shadow(s)
    return s


def set_slide_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_footer_bar(slide, page_num, total=None):
    """Dezenter Footer in Valkeen-Navy links, Seitenzahl rechts."""
    # Divider
    d = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(0.5), Inches(7.08),
                                   Inches(12.83), Inches(7.08))
    d.line.color.rgb = GREY_20
    d.line.width = Pt(0.5)
    # Left footer text
    add_text(slide, 0.5, 7.12, 10.0, 0.28,
             "Valkeen GmbH  ·  MS Project ⇄ Tempus Integration",
             size=8.5, color=GREY_40, align=PP_ALIGN.LEFT)
    # Page num
    add_text(slide, 11.8, 7.12, 1.03, 0.28,
             f"{page_num:02d}" + (f" / {total:02d}" if total else ""),
             size=8.5, color=NAVY, align=PP_ALIGN.RIGHT, bold=True)


def add_title_block(slide, title_main, accent=None, *, subtitle=None):
    """
    Valkeen-Titelzeile: schmaler Akzentbalken + Titel (Calibri Light, Navy).
    """
    # Akzentbalken
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), Inches(0.45),
                                 Inches(0.14), Inches(0.60))
    _fill(bar, GREEN)
    _no_line(bar)
    _kill_shadow(bar)

    # Titel (zwei Runs: Hauptteil + Akzent-Teil)
    tb = slide.shapes.add_textbox(Inches(0.80), Inches(0.40),
                                  Inches(12.0), Inches(0.75))
    tf = tb.text_frame
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = title_main
    r1.font.name = TITLE_FONT
    r1.font.size = Pt(26)
    r1.font.bold = True
    r1.font.color.rgb = NAVY
    if accent:
        r2 = p.add_run()
        r2.text = accent
        r2.font.name = TITLE_FONT
        r2.font.size = Pt(26)
        r2.font.bold = True
        r2.font.color.rgb = SLATE

    if subtitle:
        add_text(slide, 0.80, 1.05, 12.0, 0.35, subtitle,
                 size=12, italic=True, color=TEXT_MUTED, align=PP_ALIGN.LEFT,
                 font=BODY_FONT)


# ===========================================================================
# SLIDE 1 – Technical Boundary: Schedule vs. Allocation (1:1 Nachbau)
# ===========================================================================
def slide_technical_boundary(prs, page_num, total=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_block(
        slide,
        "The Technical Boundary: ",
        "Schedule vs. Allocation",
    )

    # Golden-Rule-Box (oben) – hell mit Navy-Rahmen links (Key-Icon-Feeling)
    gr = add_rect(slide, 0.50, 1.45, 6.30, 0.95,
                  fill=GREY_10, line=GREY_20, line_w=0.75, rounded=True,
                  radius=0.05)
    # "Key"-Kreis (Ikonersatz)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(0.68), Inches(1.66),
                                  Inches(0.52), Inches(0.52))
    _fill(circ, NAVY)
    _no_line(circ)
    _kill_shadow(circ)
    _text(circ, "🔑", size=14, color=WHITE, bold=True)

    add_text(slide, 1.34, 1.52, 5.40, 0.85,
             "The Golden Rule:  The destination of your\n"
             "data dictates your integration tool.",
             size=11, color=TEXT_DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE)

    # Source-Block links unten
    src = add_rect(slide, 0.50, 3.55, 2.35, 1.15, fill=NAVY, line=None,
                   rounded=False)
    _text(src, "MS Project Data", size=15, bold=True, color=WHITE,
          font=TITLE_FONT)

    # =======================================================================
    # Oberer Pfad – Allocations (Orange, 3 Pills)
    # =======================================================================
    top_y = 3.00    # y-Mitte des oberen Pfeils
    arrow_h = 1.25  # Dicke
    # Breiter Pfeil von (x=2.90) bis (x=10.45) – führt in den Zielblock
    add_right_arrow_block(slide, 2.90, top_y, 7.70, arrow_h,
                          fill=ORANGE, head_ratio=0.08)
    # drei Pills (weiß mit orangem Rand, wie Screenshot)
    pill_y = top_y + 0.27
    pill_h = 0.70
    # Pill 1: Project Online Sync
    add_pill(slide, 3.15, pill_y, 2.00, pill_h, "Project Online\nSync",
             fill=WHITE, text_color=NAVY, size=10.5, line=ORANGE)
    # Pill 2: MSP Add-In
    add_pill(slide, 5.40, pill_y, 1.80, pill_h, "MSP Add-In",
             fill=WHITE, text_color=NAVY, size=11, line=ORANGE)
    # Pill 3: MPP Import
    add_pill(slide, 7.45, pill_y, 1.80, pill_h, "MPP Import",
             fill=WHITE, text_color=NAVY, size=11, line=ORANGE)

    # Zielblock oben: Allocations Grid
    tgt_top = add_rect(slide, 10.85, top_y + 0.08, 1.95, 1.10,
                       fill=ORANGE, line=None, rounded=False)
    _text(tgt_top, "Allocations\nGrid", size=15, bold=True, color=WHITE,
          font=TITLE_FONT)
    # Beschreibung darunter
    add_text(slide, 10.75, top_y + 1.22, 2.10, 0.95,
             "High-level planned and actual\n"
             "resource demand. Supported by\n"
             "all synchronisation methods.",
             size=8.5, color=TEXT_MUTED, align=PP_ALIGN.LEFT)

    # =======================================================================
    # Unterer Pfad – Full Schedule (Cyan-Blau, nur MPP Import)
    # =======================================================================
    bot_y = 5.10
    add_right_arrow_block(slide, 2.90, bot_y, 7.70, arrow_h,
                          fill=SLATE, head_ratio=0.08)
    # Eine Pill in der Mitte
    add_pill(slide, 5.40, bot_y + 0.27, 1.80, pill_h, "MPP Import",
             fill=WHITE, text_color=NAVY, size=11, line=SLATE)

    # Zielblock unten: Full Schedule
    tgt_bot = add_rect(slide, 10.85, bot_y + 0.08, 1.95, 1.10,
                       fill=SLATE, line=None, rounded=False)
    _text(tgt_bot, "Full Schedule", size=15, bold=True, color=WHITE,
          font=TITLE_FONT)
    add_text(slide, 10.75, bot_y + 1.22, 2.10, 1.00,
             "Granular task logic, dependencies,\n"
             "task types, and dates.\n"
             "Supported only by MPP Import.",
             size=8.5, color=TEXT_MUTED, align=PP_ALIGN.LEFT)

    # Abzweig-Konnektoren vom Source-Block
    # oberer Abzweig
    c1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(2.85), Inches(3.75),
                                    Inches(2.95), Inches(3.20))
    c1.line.color.rgb = ORANGE
    c1.line.width = Pt(4)
    # unterer Abzweig
    c2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(2.85), Inches(4.50),
                                    Inches(2.95), Inches(5.30))
    c2.line.color.rgb = SLATE
    c2.line.width = Pt(4)

    add_footer_bar(slide, page_num, total)


# ===========================================================================
# SLIDE 2 – Integration Options: At a Glance (1:1 Nachbau)
# ===========================================================================
def slide_integration_options(prs, page_num, total=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_block(slide, "Integration Options: ", "At a Glance")

    # Tabellen-Geometrie
    left = 0.50
    top = 1.55
    col_label_w = 1.50
    col_w = 3.78
    total_w = col_label_w + 3 * col_w

    # Row-Höhen
    h_header = 0.80
    h_target = 0.95
    h_pros   = 2.10
    h_cons   = 1.60

    def cell(x, y, w, h, *, fill=WHITE, line=GREY_20, rounded=False):
        return add_rect(slide, x, y, w, h, fill=fill, line=line,
                        line_w=0.75, rounded=rounded)

    # ----- Header-Row -----
    cell(left, top, col_label_w, h_header, fill=WHITE, line=None)  # leer
    headers = [
        ("Project Online Sync", NAVY),
        ("MSP Add-In",         NAVY),
        ("MPP Import",         NAVY),
    ]
    for i, (txt, col) in enumerate(headers):
        x = left + col_label_w + i * col_w
        c = cell(x, top, col_w, h_header, fill=WHITE, line=None)
        _text(c, txt, size=17, bold=True, color=col, font=TITLE_FONT)
    # Dezente Trennlinie unter Headern
    sep = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     Inches(left + col_label_w),
                                     Inches(top + h_header),
                                     Inches(left + total_w),
                                     Inches(top + h_header))
    sep.line.color.rgb = GREY_20
    sep.line.width = Pt(0.75)

    # ----- Row: Data Target -----
    y = top + h_header
    lbl = cell(left, y, col_label_w, h_target, fill=GREY_10, line=GREY_20)
    _text(lbl, "Data\nTarget", size=12, bold=True, color=NAVY,
          font=BODY_FONT)

    for i in range(3):
        cx = left + col_label_w + i * col_w
        cell(cx, y, col_w, h_target, fill=WHITE, line=GREY_20)

    # Pills – Data Target
    #   Spalten 1 & 2: Allocations ONLY (orange)
    for col_idx in (0, 1):
        cx = left + col_label_w + col_idx * col_w
        add_pill(slide, cx + (col_w - 2.00) / 2, y + 0.30, 2.00, 0.38,
                 "Allocations ONLY", fill=ORANGE, size=10.5)
    #   Spalte 3: Full Schedule + Allocations
    cx = left + col_label_w + 2 * col_w
    # zwei Pills nebeneinander mittig
    pill1_w = 1.55
    pill2_w = 1.45
    gap = 0.12
    group_w = pill1_w + pill2_w + gap
    start_x = cx + (col_w - group_w) / 2
    add_pill(slide, start_x, y + 0.30, pill1_w, 0.38, "Full Schedule",
             fill=SLATE, size=10.5)
    add_pill(slide, start_x + pill1_w + gap, y + 0.30, pill2_w, 0.38,
             "Allocations", fill=ORANGE, size=10.5)

    # ----- Row: Pros -----
    y = top + h_header + h_target
    lbl = cell(left, y, col_label_w, h_pros, fill=GREY_10, line=GREY_20)
    _text(lbl, "Pros", size=12, bold=True, color=NAVY, font=BODY_FONT)

    pros_content = [
        [
            "Enables automated, scheduled synchronisation",
            "Admin-led governance ensures systemic control",
            "Centralised process",
        ],
        [
            "Complete Project Manager control",
            "Features a pre-sync data validation and error preview interface",
            "Enables advanced WBS mapping",
        ],
        [
            "The only tool that imports full schedules (tasks, dependencies, dates, and durations)",
            "Excellent for initial system migrations",
        ],
    ]
    for i, items in enumerate(pros_content):
        cx = left + col_label_w + i * col_w
        cell(cx, y, col_w, h_pros, fill=WHITE, line=GREY_20)
        tb = slide.shapes.add_textbox(Inches(cx + 0.22), Inches(y + 0.14),
                                      Inches(col_w - 0.30),
                                      Inches(h_pros - 0.20))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        for k, line in enumerate(items):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.level = 0
            p.space_after = Pt(4)
            r1 = p.add_run()
            r1.text = "✓  "
            r1.font.size = Pt(11)
            r1.font.bold = True
            r1.font.color.rgb = GREEN
            r1.font.name = BODY_FONT
            r2 = p.add_run()
            r2.text = line
            r2.font.size = Pt(10.5)
            r2.font.color.rgb = TEXT_DARK
            r2.font.name = BODY_FONT

    # ----- Row: Cons -----
    y = top + h_header + h_target + h_pros
    lbl = cell(left, y, col_label_w, h_cons, fill=GREY_10, line=GREY_20)
    _text(lbl, "Cons", size=12, bold=True, color=NAVY, font=BODY_FONT)

    cons_content = [
        [
            "Physically incapable of importing schedule data",
            "Requires root URL and global API credential setup",
        ],
        [
            "Requires manual execution per project",
            "Requires individual user API token generation",
        ],
        [
            "A manual, file-based process requiring drag-and-drop execution",
            "Schedule imports automatically overwrite existing schedule entities",
        ],
    ]
    for i, items in enumerate(cons_content):
        cx = left + col_label_w + i * col_w
        cell(cx, y, col_w, h_cons, fill=WHITE, line=GREY_20)
        tb = slide.shapes.add_textbox(Inches(cx + 0.22), Inches(y + 0.14),
                                      Inches(col_w - 0.30),
                                      Inches(h_cons - 0.20))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        for k, line in enumerate(items):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(4)
            r1 = p.add_run()
            r1.text = "!  "
            r1.font.size = Pt(11)
            r1.font.bold = True
            r1.font.color.rgb = ORANGE
            r1.font.name = BODY_FONT
            r2 = p.add_run()
            r2.text = line
            r2.font.size = Pt(10.5)
            r2.font.color.rgb = TEXT_DARK
            r2.font.name = BODY_FONT

    add_footer_bar(slide, page_num, total)


# ===========================================================================
# SLIDE 3 – Selecting Your Integration Strategy (1:1 Nachbau)
# ===========================================================================
def _icon_converge(slide, cx, cy):
    """4 Pfeile, die auf einen zentralen Kreis zeigen."""
    # zentraler dunkler Kreis
    r = 0.28
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(cx - r), Inches(cy - r),
                                  Inches(r * 2), Inches(r * 2))
    _fill(circ, NAVY)
    _no_line(circ)
    _kill_shadow(circ)

    # vier Pfeile (rechts/links/oben/unten), die nach innen zeigen
    arrow_len = 0.55
    arrow_w   = 0.28
    # von links nach rechts (zeigt nach rechts, landet am Kreis)
    def add_arrow(rot, dx, dy):
        s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(cx + dx), Inches(cy + dy),
                                   Inches(arrow_len), Inches(arrow_w))
        s.rotation = rot
        _fill(s, NAVY)
        _no_line(s)
        _kill_shadow(s)
        try:
            s.adjustments[0] = 0.55
            s.adjustments[1] = 0.30
        except Exception:
            pass
        return s

    # Links -> Mitte  (unrotated points right)
    add_arrow(0,   -0.95, -arrow_w / 2)
    # Rechts -> Mitte (rotate 180)
    add_arrow(180,  0.40, -arrow_w / 2)
    # Oben -> Mitte (rotate 90)
    add_arrow(90,  -arrow_w / 2 + 0.14, -0.95 - (arrow_len - arrow_w) / 2)
    # Unten -> Mitte (rotate 270)
    add_arrow(270, -arrow_w / 2 + 0.14,  0.40 - (arrow_len - arrow_w) / 2)


def _icon_network(slide, cx, cy):
    """Hub-and-Spoke: zentraler großer Knoten + 6 kleine Knoten drumherum."""
    import math
    R = 0.90  # Radius
    # zentraler Kreis
    big_r = 0.32
    big = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(cx - big_r), Inches(cy - big_r),
                                 Inches(big_r * 2), Inches(big_r * 2))
    _fill(big, NAVY)
    _no_line(big)
    _kill_shadow(big)

    # 6 kleine Knoten im Kreis
    small_r = 0.22
    for k in range(6):
        ang = math.radians(60 * k - 90)  # oben beginnen
        sx = cx + R * math.cos(ang)
        sy = cy + R * math.sin(ang)
        # Linie vom Zentrum zum Knoten (unter den Kreis)
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        Inches(cx), Inches(cy),
                                        Inches(sx), Inches(sy))
        ln.line.color.rgb = NAVY
        ln.line.width = Pt(1.25)
        # kleiner Knoten (Kopf)
        sc = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(sx - small_r),
                                    Inches(sy - small_r),
                                    Inches(small_r * 2),
                                    Inches(small_r * 2))
        _fill(sc, NAVY)
        _no_line(sc)
        _kill_shadow(sc)
        # "Körper" als kleiner Halbkreis darunter
        body = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(sx - small_r * 1.25),
                                      Inches(sy + small_r * 0.35),
                                      Inches(small_r * 2.5),
                                      Inches(small_r * 1.2))
        _fill(body, NAVY)
        _no_line(body)
        _kill_shadow(body)


def _icon_disk(slide, cx, cy):
    """HDD mit Pfeil – vereinfachtes Disk-Migrations-Icon."""
    # Scheiben-Stapel (3 Ovale übereinander)
    w = 1.10
    h_disk = 0.22
    spacing = 0.20
    for i in range(3):
        dy = cy - 0.55 + i * spacing
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(cx - w / 2), Inches(dy),
                                   Inches(w), Inches(h_disk))
        _fill(d, NAVY)
        _no_line(d)
        _kill_shadow(d)
    # Verbindungsrechteck (vertikal, dunkel)
    body = add_rect(slide, cx - w / 2, cy - 0.44,
                    w, 0.20 * 2 + 0.05,
                    fill=NAVY, line=None)
    # runder Pfeil nach rechts (Migration)
    arr = slide.shapes.add_shape(MSO_SHAPE.CURVED_RIGHT_ARROW,
                                 Inches(cx + 0.25), Inches(cy - 0.85),
                                 Inches(0.80), Inches(1.50))
    _fill(arr, NAVY)
    _no_line(arr)
    _kill_shadow(arr)


def slide_selecting_strategy(prs, page_num, total=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_title_block(slide, "Selecting Your ", "Integration Strategy")

    cards = [
        {
            "icon": "converge",
            "situation": "If your operational model is…",
            "body": "Centralised planning requiring automated, hands-off "
                    "synchronisation across the enterprise.",
            "arrow_text": "Then deploy",
            "tool": "Project\nOnline Sync",
            "accent": ORANGE,
        },
        {
            "icon": "network",
            "situation": "If your operational model is…",
            "body": "Decentralised planning where individual Project Managers "
                    "must control, preview, and validate their own data "
                    "before submission.",
            "arrow_text": "Then deploy the",
            "tool": "MSP\nAdd-In",
            "accent": ORANGE,
        },
        {
            "icon": "disk",
            "situation": "If your operational model requires…",
            "body": "A one-time system migration, or an ongoing necessity to "
                    "view full task logic, dependencies, and granular "
                    "project schedules in Tempus.",
            "arrow_text": "Then utilise",
            "tool": "MPP\nImport",
            "accent": SLATE,
        },
    ]

    card_y = 1.55
    card_w = 4.08
    card_h = 5.35
    gap    = 0.16
    start_x = 0.50

    for i, c in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rect(slide, x, card_y, card_w, card_h,
                 fill=WHITE, line=GREY_20, line_w=0.75, rounded=False)

        # Icon-Bereich
        icon_cx = x + card_w / 2
        icon_cy = card_y + 1.00
        if c["icon"] == "converge":
            _icon_converge(slide, icon_cx, icon_cy)
        elif c["icon"] == "network":
            _icon_network(slide, icon_cx, icon_cy)
        elif c["icon"] == "disk":
            _icon_disk(slide, icon_cx, icon_cy)

        # Situation
        add_text(slide, x + 0.30, card_y + 2.15, card_w - 0.60, 0.40,
                 c["situation"], size=12, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, font=BODY_FONT)

        # Body
        add_text(slide, x + 0.30, card_y + 2.55, card_w - 0.60, 1.30,
                 c["body"], size=11, color=TEXT_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
                 font=BODY_FONT, line_spacing=1.2)

        # Abwärtspfeil
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     Inches(icon_cx - 0.24),
                                     Inches(card_y + 3.90),
                                     Inches(0.48), Inches(0.45))
        _fill(arr, NAVY)
        _no_line(arr)
        _kill_shadow(arr)

        # Then deploy / utilise
        add_text(slide, x + 0.30, card_y + 4.40, card_w - 0.60, 0.30,
                 c["arrow_text"], size=11, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, font=BODY_FONT)

        # Tool-Name (groß) – Akzentfarbe
        add_text(slide, x + 0.30, card_y + 4.70, card_w - 0.60, 0.65,
                 c["tool"], size=22, bold=True, color=c["accent"],
                 align=PP_ALIGN.CENTER, font=TITLE_FONT,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.0)

    add_footer_bar(slide, page_num, total)


# ===========================================================================
# (Optional) Alternativ-Varianten – werden ans Deck angehängt
# ===========================================================================
def slide_boundary_variant_swimlane(prs, page_num, total=None):
    """Variante B zu Slide 1 – als Lane-Diagramm."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_block(slide, "Technische Abgrenzung · ", "Alternativdarstellung",
                    subtitle="Zwei Ziel-Datendomänen – jeweils mit den "
                             "unterstützten Integrationswerkzeugen.")

    # Source
    src = add_rect(slide, 0.50, 2.40, 1.95, 3.10, fill=NAVY, line=None)
    _text(src, "MS Project\nData", size=16, bold=True, color=WHITE,
          font=TITLE_FONT)

    # Lane 1 (Allocations)
    y1 = 1.70
    lane1 = add_rect(slide, 2.70, y1, 10.10, 2.15,
                     fill=ORANGE_LT, line=ORANGE, line_w=1.0, rounded=True,
                     radius=0.04)
    add_text(slide, 2.95, y1 + 0.10, 4.0, 0.40, "Allocations Grid",
             size=16, bold=True, color=ORANGE, font=TITLE_FONT)
    add_text(slide, 2.95, y1 + 0.50, 4.0, 0.50,
             "High-level Planung & Ressourcenbedarf",
             size=10, color=TEXT_MUTED, font=BODY_FONT)
    # Drei Chevrons
    cw, ch = 1.95, 0.90
    cy = y1 + 1.08
    cx0 = 6.95
    for i, label in enumerate(["Project Online Sync", "MSP Add-In", "MPP Import"]):
        s = slide.shapes.add_shape(MSO_SHAPE.PENTAGON,
                                   Inches(cx0 + i * (cw + 0.03)),
                                   Inches(cy), Inches(cw), Inches(ch))
        _fill(s, ORANGE)
        _no_line(s)
        _kill_shadow(s)
        _text(s, label, size=11, bold=True, color=WHITE, font=BODY_FONT)

    # Lane 2 (Full Schedule)
    y2 = 4.20
    add_rect(slide, 2.70, y2, 10.10, 2.15,
             fill=CYAN_LIGHT, line=SLATE, line_w=1.0, rounded=True,
             radius=0.04)
    add_text(slide, 2.95, y2 + 0.10, 4.0, 0.40, "Full Schedule",
             size=16, bold=True, color=SLATE, font=TITLE_FONT)
    add_text(slide, 2.95, y2 + 0.50, 4.0, 0.50,
             "Tasks, Dependencies, Dauern, Termine",
             size=10, color=TEXT_MUTED, font=BODY_FONT)
    cy = y2 + 1.08
    s = slide.shapes.add_shape(MSO_SHAPE.PENTAGON,
                               Inches(cx0 + (cw + 0.03)),
                               Inches(cy), Inches(cw), Inches(ch))
    _fill(s, SLATE)
    _no_line(s)
    _kill_shadow(s)
    _text(s, "MPP Import", size=11, bold=True, color=WHITE, font=BODY_FONT)

    # Konnektoren von Source zu Lanes
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(2.45), Inches(3.20),
                                    Inches(2.70), Inches(2.78))
    ln.line.color.rgb = ORANGE
    ln.line.width = Pt(3)
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(2.45), Inches(4.70),
                                    Inches(2.70), Inches(5.28))
    ln.line.color.rgb = SLATE
    ln.line.width = Pt(3)

    add_footer_bar(slide, page_num, total)


def slide_options_variant_cards(prs, page_num, total=None):
    """Variante B zu Slide 2 – 3 Tile-Cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_block(slide, "Integrationsoptionen · ", "Alternativdarstellung",
                    subtitle="Drei Wege, um MS-Project-Daten nach Tempus zu bringen.")

    options = [
        {
            "title": "Project Online Sync",
            "badge": [("Allocations ONLY", ORANGE)],
            "pros": [
                "Automatisierte, geplante Synchronisation",
                "Admin-gesteuerte Governance",
                "Zentralisierter Prozess",
            ],
            "cons": [
                "Keine Schedule-Daten importierbar",
                "Root-URL & globale API-Credentials nötig",
            ],
            "accent": NAVY,
        },
        {
            "title": "MSP Add-In",
            "badge": [("Allocations ONLY", ORANGE)],
            "pros": [
                "Volle Kontrolle beim Project Manager",
                "Datenvalidierung & Vorschau vor Sync",
                "Erweitertes WBS-Mapping",
            ],
            "cons": [
                "Manuelle Ausführung pro Projekt",
                "Individuelle API-Token je Nutzer",
            ],
            "accent": NAVY,
        },
        {
            "title": "MPP Import",
            "badge": [("Full Schedule", SLATE), ("Allocations", ORANGE)],
            "pros": [
                "Einziges Tool für vollständige Schedules\n"
                "(Tasks, Dependencies, Dauer, Termine)",
                "Ideal für initiale Systemmigrationen",
            ],
            "cons": [
                "Manueller Drag-&-Drop-Prozess (dateibasiert)",
                "Überschreibt bestehende Schedule-Daten",
            ],
            "accent": NAVY,
        },
    ]

    card_y = 1.55
    card_w = 4.05
    card_h = 5.30
    gap = 0.15
    start_x = 0.55

    for i, opt in enumerate(options):
        x = start_x + i * (card_w + gap)
        add_rect(slide, x, card_y, card_w, card_h,
                 fill=WHITE, line=GREY_20, line_w=0.75)
        head = add_rect(slide, x, card_y, card_w, 0.80,
                        fill=opt["accent"], line=None)
        _text(head, opt["title"], size=16, bold=True, color=WHITE,
              font=TITLE_FONT)

        bx = x + 0.28
        by = card_y + 0.95
        for label, col in opt["badge"]:
            bw = 1.55 if len(label) < 16 else 1.90
            add_pill(slide, bx, by, bw, 0.38, label, fill=col, size=10.5)
            bx += bw + 0.10

        py = card_y + 1.55
        add_text(slide, x + 0.28, py, card_w - 0.56, 0.30, "Vorteile",
                 size=11, bold=True, color=GREEN, font=BODY_FONT)
        tb = slide.shapes.add_textbox(Inches(x + 0.28), Inches(py + 0.32),
                                      Inches(card_w - 0.56), Inches(1.95))
        tf = tb.text_frame
        tf.word_wrap = True
        for k, line in enumerate(opt["pros"]):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.space_after = Pt(4)
            r1 = p.add_run()
            r1.text = "✓  "
            r1.font.size = Pt(11)
            r1.font.bold = True
            r1.font.color.rgb = GREEN
            r1.font.name = BODY_FONT
            r2 = p.add_run()
            r2.text = line
            r2.font.size = Pt(10.5)
            r2.font.color.rgb = TEXT_DARK
            r2.font.name = BODY_FONT

        cy = card_y + 3.60
        add_text(slide, x + 0.28, cy, card_w - 0.56, 0.30, "Nachteile",
                 size=11, bold=True, color=ORANGE, font=BODY_FONT)
        tb = slide.shapes.add_textbox(Inches(x + 0.28), Inches(cy + 0.32),
                                      Inches(card_w - 0.56), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        for k, line in enumerate(opt["cons"]):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.space_after = Pt(4)
            r1 = p.add_run()
            r1.text = "!  "
            r1.font.size = Pt(11)
            r1.font.bold = True
            r1.font.color.rgb = ORANGE
            r1.font.name = BODY_FONT
            r2 = p.add_run()
            r2.text = line
            r2.font.size = Pt(10.5)
            r2.font.color.rgb = TEXT_DARK
            r2.font.name = BODY_FONT

    add_footer_bar(slide, page_num, total)


def slide_strategy_variant_matrix(prs, page_num, total=None):
    """Variante B zu Slide 3 – Entscheidungs-Matrix."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_block(slide, "Entscheidungsmatrix · ",
                    "Welches Tool passt wann?",
                    subtitle="Operativer Rahmen × Merkmale × Empfohlenes Tool.")

    top = 1.60
    left = 0.50
    col_w = [4.30, 5.40, 2.93]
    row_h_head = 0.55

    def cell(x, y, w, h, *, fill, line=GREY_20, text="", size=12, bold=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, font=BODY_FONT):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y),
                                   Inches(w), Inches(h))
        _fill(s, fill)
        _line(s, line, 0.75)
        _kill_shadow(s)
        if text:
            _text(s, text, size=size, bold=bold, color=color, align=align,
                  font=font)
        return s

    cell(left, top, col_w[0], row_h_head, fill=NAVY, line=NAVY,
         text="Operativer Rahmen", size=13, bold=True, color=WHITE,
         font=TITLE_FONT, align=PP_ALIGN.LEFT)
    cell(left + col_w[0], top, col_w[1], row_h_head, fill=NAVY, line=NAVY,
         text="Merkmale", size=13, bold=True, color=WHITE, font=TITLE_FONT,
         align=PP_ALIGN.LEFT)
    cell(left + col_w[0] + col_w[1], top, col_w[2], row_h_head,
         fill=NAVY, line=NAVY,
         text="Empfohlenes Tool", size=13, bold=True, color=WHITE,
         font=TITLE_FONT, align=PP_ALIGN.CENTER)

    rows = [
        {
            "op": "Zentrale Planung",
            "desc": "Automatisierte, hands-off Synchronisation unternehmensweit. "
                    "Governance liegt beim Admin.",
            "tool": "Project\nOnline Sync",
            "accent": ORANGE,
        },
        {
            "op": "Dezentrale Planung",
            "desc": "Einzelne Project Manager steuern, prüfen und validieren "
                    "ihre Daten vor dem Transfer selbst.",
            "tool": "MSP\nAdd-In",
            "accent": ORANGE,
        },
        {
            "op": "Migration / Volle Tiefe",
            "desc": "Einmalige Systemmigration oder laufender Bedarf an vollständiger "
                    "Task-Logik, Dependencies und Terminen in Tempus.",
            "tool": "MPP\nImport",
            "accent": SLATE,
        },
    ]
    y = top + row_h_head
    row_h = 1.45
    for i, r in enumerate(rows):
        bg = GREY_10 if i % 2 == 0 else WHITE
        cell(left, y, col_w[0], row_h, fill=bg, text=r["op"],
             size=14, bold=True, color=NAVY, font=TITLE_FONT,
             align=PP_ALIGN.LEFT)
        cell(left + col_w[0], y, col_w[1], row_h, fill=bg, text=r["desc"],
             size=11, color=TEXT_MUTED, align=PP_ALIGN.LEFT)
        cell(left + col_w[0] + col_w[1], y, col_w[2], row_h,
             fill=r["accent"], line=r["accent"], text=r["tool"],
             size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             font=TITLE_FONT)
        y += row_h

    add_text(slide, 0.50, 6.60, 12.3, 0.35,
             "Auswahlregel: Zielsystem (Allocations vs. Full Schedule) × "
             "operatives Modell → passendes Integrationswerkzeug.",
             size=10.5, italic=True, color=TEXT_MUTED,
             align=PP_ALIGN.LEFT)

    add_footer_bar(slide, page_num, total)


# ===========================================================================
# Titel- und Divider-Folien
# ===========================================================================
def slide_title_deck(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    # Akzentbalken
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.8), Inches(2.95),
                                 Inches(0.14), Inches(1.30))
    _fill(bar, GREEN)
    _no_line(bar)
    _kill_shadow(bar)

    add_text(slide, 1.1, 2.60, 11.5, 0.55,
             "TEMPUS · MS PROJECT INTEGRATION",
             size=14, bold=True, color=GREEN, align=PP_ALIGN.LEFT,
             font=BODY_FONT)
    add_text(slide, 1.1, 3.00, 11.5, 1.20,
             "Integrations-\noptionen",
             size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
             font=TITLE_FONT, line_spacing=1.0)
    add_text(slide, 1.1, 4.50, 11.5, 0.60,
             "Editierbare Folienvarianten im Valkeen-Design",
             size=16, color=RGBColor(0xCF, 0xD4, 0xDB), align=PP_ALIGN.LEFT,
             font=BODY_FONT)

    # Footer
    add_text(slide, 0.8, 6.95, 11.5, 0.30,
             "Valkeen GmbH  ·  www.valkeen.com",
             size=10, color=GREY_40, align=PP_ALIGN.LEFT, font=BODY_FONT)


def slide_agenda(prs, page_num, total=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_block(slide, "Inhalt")

    items = [
        ("01", "Technical Boundary",
         "Schedule vs. Allocation – Datenziel steuert das Werkzeug."),
        ("02", "Integration Options",
         "At a Glance – Pros & Cons der drei Werkzeuge im Vergleich."),
        ("03", "Selecting Your Strategy",
         "Operatives Modell → passendes Integrationswerkzeug."),
        ("04", "Alternativdarstellungen",
         "Swimlane · Tile-Cards · Entscheidungsmatrix."),
    ]
    y = 1.75
    for num, title, sub in items:
        # Nummer
        nbar = add_rect(slide, 0.55, y, 0.65, 0.95,
                        fill=NAVY, line=None, rounded=False)
        _text(nbar, num, size=20, bold=True, color=WHITE, font=TITLE_FONT)
        # Grüner Strich
        gl = add_rect(slide, 1.25, y + 0.08, 0.08, 0.80,
                      fill=GREEN, line=None)
        _kill_shadow(gl)
        # Titel
        add_text(slide, 1.50, y + 0.05, 11.0, 0.45, title,
                 size=17, bold=True, color=NAVY, font=TITLE_FONT)
        add_text(slide, 1.50, y + 0.50, 11.0, 0.45, sub,
                 size=11, color=TEXT_MUTED, font=BODY_FONT)
        y += 1.15

    add_footer_bar(slide, page_num, total)


# ===========================================================================
# Build presentations
# ===========================================================================
def build_standalone(out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = 8  # 1 title + 1 agenda + 3 main + 3 variants
    slide_title_deck(prs)
    slide_agenda(prs, 2, total)
    slide_technical_boundary(prs, 3, total)
    slide_integration_options(prs, 4, total)
    slide_selecting_strategy(prs, 5, total)
    slide_boundary_variant_swimlane(prs, 6, total)
    slide_options_variant_cards(prs, 7, total)
    slide_strategy_variant_matrix(prs, 8, total)

    prs.save(out_path)
    print(f"OK · standalone -> {out_path}")


def build_on_template(template_path, out_path):
    """Lädt die Knauf-Vorlage und hängt unsere neuen Folien an."""
    prs = Presentation(template_path)
    start_idx = len(prs.slides) + 1
    total = start_idx + 5

    slide_technical_boundary(prs, start_idx, total)
    slide_integration_options(prs, start_idx + 1, total)
    slide_selecting_strategy(prs, start_idx + 2, total)
    slide_boundary_variant_swimlane(prs, start_idx + 3, total)
    slide_options_variant_cards(prs, start_idx + 4, total)
    slide_strategy_variant_matrix(prs, start_idx + 5, total)

    prs.save(out_path)
    print(f"OK · on template -> {out_path}")


if __name__ == "__main__":
    out_dir = "PowerPoint-Exports"
    os.makedirs(out_dir, exist_ok=True)

    build_standalone(os.path.join(out_dir, "Valkeen_Integration-Optionen.pptx"))

    template = "/Users/manumanera/Downloads/Entscheidungsvorlage Knauf.pptx"
    if os.path.exists(template):
        build_on_template(
            template,
            os.path.join(
                out_dir,
                "Entscheidungsvorlage_Knauf_mit_Integrationsfolien.pptx",
            ),
        )
    else:
        print(f"WARN: Template nicht gefunden: {template}")
