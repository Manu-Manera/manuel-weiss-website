"""
Erstellt editierbare PowerPoint-Folien (100 % Vektor/Shapes, keine Bilder)
basierend auf den drei Screenshots zum Thema
"MS Project / Tempus Integrationsoptionen".

Für jedes Originaldiagramm werden zwei Varianten ausgegeben:
  - Variante A: möglichst 1:1-Nachbau
  - Variante B: Corporate-Alternative / SmartArt-ähnliche Darstellung

Ausgabe:
  PowerPoint-Exports/Integration-Optionen_editierbar.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Farbpalette (an Screenshots angelehnt)
# ---------------------------------------------------------------------------
CLR_DARK        = RGBColor(0x1F, 0x2A, 0x3A)   # sehr dunkles Blau/Anthrazit
CLR_ORANGE      = RGBColor(0xF2, 0x9B, 0x2C)   # Allocations-Orange
CLR_ORANGE_SOFT = RGBColor(0xFF, 0xE9, 0xCC)
CLR_BLUE        = RGBColor(0x1F, 0x4E, 0x9C)   # Full-Schedule-Blau
CLR_BLUE_SOFT   = RGBColor(0xD8, 0xE4, 0xF4)
CLR_BG          = RGBColor(0xF5, 0xF6, 0xF8)
CLR_CARD        = RGBColor(0xFF, 0xFF, 0xFF)
CLR_BORDER      = RGBColor(0xD9, 0xDD, 0xE3)
CLR_TEXT        = RGBColor(0x1F, 0x2A, 0x3A)
CLR_MUTED       = RGBColor(0x6B, 0x72, 0x80)
CLR_GREEN       = RGBColor(0x2E, 0x7D, 0x32)
CLR_RED         = RGBColor(0xC6, 0x28, 0x28)

# 16:9 in Inches
SLIDE_W, SLIDE_H = 13.333, 7.5


# ---------------------------------------------------------------------------
# Hilfsfunktionen
# ---------------------------------------------------------------------------
def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_fill(shape):
    shape.fill.background()


def set_line(shape, color, width_pt=1.0):
    ln = shape.line
    ln.color.rgb = color
    ln.width = Pt(width_pt)


def no_line(shape):
    shape.line.fill.background()


def set_text(shape, text, *, size=14, bold=False, color=CLR_TEXT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = anchor
    tf.clear()
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_rect(slide, x, y, w, h, *, fill=CLR_CARD, line=CLR_BORDER, line_w=1.0,
             rounded=True, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    if rounded:
        # kleineren Radius setzen
        s.adjustments[0] = 0.12
    if fill is None:
        no_fill(s)
    else:
        set_fill(s, fill)
    if line is None:
        no_line(s)
    else:
        set_line(s, line, line_w)
    if not shadow:
        # Schatten deaktivieren
        sppr = s.shadow
        sppr.inherit = False
    return s


def add_pill(slide, x, y, w, h, text, *, fill=CLR_ORANGE, text_color=RGBColor(0xFF, 0xFF, 0xFF),
             size=12, bold=True, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.adjustments[0] = 0.5  # echte Pille
    set_fill(s, fill)
    if line is None:
        no_line(s)
    else:
        set_line(s, line, 0.75)
    set_text(s, text, size=size, bold=bold, color=text_color)
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=CLR_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(tb, text, size=size, bold=bold, color=color, align=align,
             anchor=anchor, font=font)
    return tb


def add_arrow(slide, x1, y1, x2, y2, *, color=CLR_ORANGE, width_pt=4.0):
    """Pfeil (Linie mit Pfeilspitze) zwischen zwei Punkten."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    ln = conn.line
    ln.color.rgb = color
    ln.width = Pt(width_pt)
    # Pfeilspitze am Ende
    lnXml = conn.line._get_or_add_ln()
    tail = etree.SubElement(lnXml, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    return conn


def add_right_arrow_block(slide, x, y, w, h, *, fill, line=None):
    """Breiter Pfeil als Block (MSO_SHAPE.RIGHT_ARROW)."""
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               Inches(x), Inches(y),
                               Inches(w), Inches(h))
    set_fill(s, fill)
    if line is None:
        no_line(s)
    else:
        set_line(s, line, 0.75)
    return s


def slide_title(slide, title_main, title_accent=None, *, y=0.35):
    """Titelzeile wie in den Screenshots."""
    if title_accent:
        add_text(slide, 0.5, y, 12.3, 0.7,
                 title_main, size=30, bold=True, color=CLR_DARK,
                 align=PP_ALIGN.LEFT)
        # Akzent-Text in zweiter Box rechts daneben - einfach inline:
        # (python-pptx unterstützt Runs – hier machen wir es über zwei Runs)
    else:
        add_text(slide, 0.5, y, 12.3, 0.7, title_main,
                 size=30, bold=True, color=CLR_DARK, align=PP_ALIGN.LEFT)


def slide_title_runs(slide, parts, *, y=0.3, x=0.5, w=12.3, h=0.8, size=30):
    """
    parts: Liste von (text, bold, color) Tupeln – wird als eine Zeile gerendert.
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for text, bold, color in parts:
        run = p.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def set_slide_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ===========================================================================
# SLIDE BUILDERS
# ===========================================================================

def build_slide_1a(prs):
    """Variante A – 1:1-Nachbau von 'The Technical Boundary: Schedule vs. Allocation'."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, CLR_BG)

    # Titel
    slide_title_runs(slide, [
        ("The Technical Boundary: ", True, CLR_DARK),
        ("Schedule vs. Allocation", True, CLR_MUTED),
    ], y=0.3, size=30)

    # Golden-Rule-Box oben
    gr = add_rect(slide, 0.5, 1.15, 6.0, 0.9, fill=CLR_CARD, line=CLR_BORDER)
    add_text(slide, 0.95, 1.22, 5.4, 0.8,
             "🔑  The Golden Rule: The destination of your data\n"
             "dictates your integration tool.",
             size=11, bold=False, color=CLR_DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE)

    # Quelle: MS Project Data
    src = add_rect(slide, 0.6, 3.5, 2.2, 1.1, fill=CLR_DARK, line=None)
    set_text(src, "MS Project Data", size=16, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))

    # Pfad 1 (oberer Strang – Allocations)
    path_top_y = 3.1
    arrow_h = 0.35
    # gelber Dick-Pfeil Hintergrund
    add_right_arrow_block(slide, 2.85, path_top_y + 0.30, 7.8, arrow_h,
                          fill=CLR_ORANGE)

    # Drei Pills auf dem oberen Strang
    add_pill(slide, 3.10, path_top_y + 0.14, 1.85, 0.70,
             "Project Online\nSync", fill=CLR_ORANGE, size=10.5, bold=True)
    add_pill(slide, 5.25, path_top_y + 0.14, 1.55, 0.70,
             "MSP Add-In", fill=CLR_ORANGE, size=11, bold=True)
    add_pill(slide, 7.10, path_top_y + 0.14, 1.55, 0.70,
             "MPP Import", fill=CLR_ORANGE, size=11, bold=True)

    # Ziel oben
    tgt_top = add_rect(slide, 10.85, path_top_y + 0.05, 1.95, 0.95,
                       fill=CLR_ORANGE, line=None)
    set_text(tgt_top, "Allocations\nGrid", size=15, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, 10.70, path_top_y + 1.05, 2.30, 1.10,
             "High-level planned and actual\n"
             "resource demand. Supported by\n"
             "all synchronisation methods.",
             size=8.5, color=CLR_MUTED, align=PP_ALIGN.LEFT)

    # Pfad 2 (unterer Strang – Full Schedule)
    path_bot_y = 5.05
    add_right_arrow_block(slide, 2.85, path_bot_y + 0.30, 7.8, arrow_h,
                          fill=CLR_BLUE)
    add_pill(slide, 5.25, path_bot_y + 0.14, 1.55, 0.70,
             "MPP Import", fill=CLR_BLUE, size=11, bold=True)

    tgt_bot = add_rect(slide, 10.85, path_bot_y + 0.05, 1.95, 0.95,
                       fill=CLR_BLUE, line=None)
    set_text(tgt_bot, "Full Schedule", size=15, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, 10.70, path_bot_y + 1.05, 2.30, 1.10,
             "Granular task logic, dependencies,\n"
             "task types, and dates.\n"
             "Supported only by MPP Import.",
             size=8.5, color=CLR_MUTED, align=PP_ALIGN.LEFT)


def build_slide_1b(prs):
    """Variante B – Corporate / SmartArt-Stil (Prozesspfeile untereinander)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CLR_CARD)

    slide_title_runs(slide, [
        ("Technische Abgrenzung – ", True, CLR_DARK),
        ("Schedule vs. Allocation", True, CLR_ORANGE),
    ], y=0.35, size=28)
    add_text(slide, 0.5, 0.95, 12.3, 0.4,
             "Das Zielsystem der Daten entscheidet über das Integrationswerkzeug.",
             size=14, color=CLR_MUTED, align=PP_ALIGN.LEFT)

    # Zwei horizontale "Swimlanes"
    # Lane 1 – Allocations (orange)
    lane_y1 = 1.9
    add_rect(slide, 0.5, lane_y1, 12.3, 1.95,
             fill=RGBColor(0xFD, 0xF4, 0xE4), line=CLR_ORANGE, line_w=1.25)
    add_text(slide, 0.8, lane_y1 + 0.1, 5.0, 0.45,
             "Allocations Grid", size=16, bold=True, color=CLR_ORANGE,
             align=PP_ALIGN.LEFT)
    add_text(slide, 0.8, lane_y1 + 0.55, 5.0, 0.5,
             "High-level Planung & Ressourcenbedarf",
             size=10.5, color=CLR_MUTED, align=PP_ALIGN.LEFT)

    # Drei Chevrons
    cx = 5.7
    cy = lane_y1 + 0.50
    cw = 1.95
    ch = 0.95
    spacing = 0.18
    for i, label in enumerate(["Project Online Sync", "MSP Add-In", "MPP Import"]):
        s = slide.shapes.add_shape(MSO_SHAPE.PENTAGON,
                                   Inches(cx + i * (cw + spacing)),
                                   Inches(cy),
                                   Inches(cw), Inches(ch))
        set_fill(s, CLR_ORANGE)
        no_line(s)
        set_text(s, label, size=11.5, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))

    # Lane 2 – Full Schedule (blau)
    lane_y2 = 4.15
    add_rect(slide, 0.5, lane_y2, 12.3, 1.95,
             fill=RGBColor(0xE8, 0xEF, 0xF9), line=CLR_BLUE, line_w=1.25)
    add_text(slide, 0.8, lane_y2 + 0.1, 5.0, 0.45,
             "Full Schedule", size=16, bold=True, color=CLR_BLUE,
             align=PP_ALIGN.LEFT)
    add_text(slide, 0.8, lane_y2 + 0.55, 5.0, 0.5,
             "Tasks, Abhängigkeiten, Termine",
             size=10.5, color=CLR_MUTED, align=PP_ALIGN.LEFT)

    # Ein Chevron in Lane 2 (mittig)
    s = slide.shapes.add_shape(MSO_SHAPE.PENTAGON,
                               Inches(cx + (cw + spacing)),
                               Inches(lane_y2 + 0.50),
                               Inches(cw), Inches(ch))
    set_fill(s, CLR_BLUE)
    no_line(s)
    set_text(s, "MPP Import", size=11.5, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))

    # Quellblock ganz links – vertikale Bracket
    src = add_rect(slide, 0.5, 2.9, 1.7, 2.25, fill=CLR_DARK, line=None)
    set_text(src, "MS Project\nData", size=15, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    # Verbindungen
    add_arrow(slide, 2.2, 3.3, 5.7, 2.85, color=CLR_ORANGE, width_pt=3)
    add_arrow(slide, 2.2, 4.8, 7.85, 4.90, color=CLR_BLUE, width_pt=3)

    # Footer-Hinweis
    add_text(slide, 0.5, 6.5, 12.3, 0.5,
             "Merksatz: Das Zielsystem der Daten bestimmt das Integrationswerkzeug.",
             size=12, bold=True, color=CLR_DARK, align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# SLIDE 2 – Integration Options: At a Glance
# ---------------------------------------------------------------------------
def build_slide_2a(prs):
    """Variante A – Tabellen-Nachbau."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CLR_BG)

    slide_title_runs(slide, [
        ("Integration Options: ", True, CLR_DARK),
        ("At a Glance", True, CLR_MUTED),
    ], y=0.3, size=28)

    # Tabellen-Geometrie
    left = 0.6
    top = 1.15
    col_label_w = 1.5
    col_w = 3.70
    row_h = [1.0, 0.95, 2.55, 2.25]  # Header, Data Target, Pros, Cons – hier 4 "Zeilen"
    # Wir machen: Header-Row, dann Data-Target-Row, Pros-Row, Cons-Row

    # Header-Row
    hdrs = ["Project Online Sync", "MSP Add-In", "MPP Import"]
    # Label-Spalte-Header leer
    add_rect(slide, left, top, col_label_w, row_h[0], fill=CLR_CARD,
             line=CLR_BORDER)
    for i, h in enumerate(hdrs):
        r = add_rect(slide,
                     left + col_label_w + i * col_w,
                     top, col_w, row_h[0],
                     fill=CLR_CARD, line=CLR_BORDER)
        set_text(r, h, size=16, bold=True, color=CLR_DARK)

    # Row-Builder
    def row(y, label, cells, label_bg=CLR_CARD):
        r = add_rect(slide, left, y, col_label_w, None or row_h[1],
                     fill=label_bg, line=CLR_BORDER, rounded=False)
        # Höhe über row_h -> außerhalb; wir lösen das mit direktem height
        return None

    # Weil die Höhe pro Zeile variiert, bauen wir Zellen einzeln:
    def cell(x, y, w, h, fill=CLR_CARD, line=CLR_BORDER, rounded=False):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
        set_fill(s, fill)
        set_line(s, line, 0.75)
        return s

    # Zeile: Data Target
    y = top + row_h[0]
    c = cell(left, y, col_label_w, row_h[1])
    set_text(c, "Data\nTarget", size=13, bold=True, color=CLR_DARK)

    # Pills pro Spalte
    for i in range(3):
        cx = left + col_label_w + i * col_w
        cell(cx, y, col_w, row_h[1])
    # Project Online Sync – Allocations ONLY (orange)
    add_pill(slide, left + col_label_w + 0.80, y + 0.28, 2.10, 0.40,
             "Allocations ONLY", fill=CLR_ORANGE, size=11)
    # MSP Add-In – Allocations ONLY
    add_pill(slide, left + col_label_w + col_w + 0.80, y + 0.28, 2.10, 0.40,
             "Allocations ONLY", fill=CLR_ORANGE, size=11)
    # MPP Import – Full Schedule + Allocations
    add_pill(slide, left + col_label_w + 2 * col_w + 0.35, y + 0.28, 1.60, 0.40,
             "Full Schedule", fill=CLR_BLUE, size=10.5)
    add_pill(slide, left + col_label_w + 2 * col_w + 2.00, y + 0.28, 1.40, 0.40,
             "Allocations", fill=CLR_ORANGE, size=10.5)

    # Zeile: Pros
    y = top + row_h[0] + row_h[1]
    c = cell(left, y, col_label_w, row_h[2])
    set_text(c, "Pros", size=13, bold=True, color=CLR_DARK)

    pros = [
        ["- Enables automated, scheduled synchronisation",
         "- Admin-led governance ensures systemic control",
         "- Centralised process"],
        ["- Complete Project Manager control",
         "- Features a pre-sync data validation\n  and error preview interface",
         "- Enables advanced WBS mapping"],
        ["- The only tool that imports full\n  schedules (tasks, dependencies,\n  dates, and durations)",
         "- Excellent for initial system migrations"],
    ]
    for i, items in enumerate(pros):
        cx = left + col_label_w + i * col_w
        cell(cx, y, col_w, row_h[2])
        tb = slide.shapes.add_textbox(Inches(cx + 0.2), Inches(y + 0.15),
                                      Inches(col_w - 0.3), Inches(row_h[2] - 0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        for k, line in enumerate(items):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = "✓ " + line.lstrip("- ").strip()
            run.font.size = Pt(10.5)
            run.font.color.rgb = CLR_GREEN
            run.font.name = "Calibri"
            p.space_after = Pt(4)

    # Zeile: Cons
    y = top + row_h[0] + row_h[1] + row_h[2]
    c = cell(left, y, col_label_w, row_h[3])
    set_text(c, "Cons", size=13, bold=True, color=CLR_DARK)

    cons = [
        ["- Physically incapable of importing\n  schedule data",
         "- Requires root URL and global API\n  credential setup"],
        ["- Requires manual execution per project",
         "- Requires individual user API token\n  generation"],
        ["- A manual, file-based process\n  requiring drag-and-drop execution",
         "- Schedule imports automatically\n  overwrite existing schedule entities"],
    ]
    for i, items in enumerate(cons):
        cx = left + col_label_w + i * col_w
        cell(cx, y, col_w, row_h[3])
        tb = slide.shapes.add_textbox(Inches(cx + 0.2), Inches(y + 0.15),
                                      Inches(col_w - 0.3), Inches(row_h[3] - 0.2))
        tf = tb.text_frame
        tf.word_wrap = True
        for k, line in enumerate(items):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = "! " + line.lstrip("- ").strip()
            run.font.size = Pt(10.5)
            run.font.color.rgb = CLR_RED
            run.font.name = "Calibri"
            p.space_after = Pt(4)


def build_slide_2b(prs):
    """Variante B – 3 Karten nebeneinander (Corporate Tile-Layout)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CLR_CARD)

    slide_title_runs(slide, [
        ("Integrationsoptionen im ", True, CLR_DARK),
        ("Überblick", True, CLR_ORANGE),
    ], y=0.35, size=28)
    add_text(slide, 0.5, 0.95, 12.3, 0.4,
             "Drei Wege, um MS-Project-Daten nach Tempus zu bringen.",
             size=14, color=CLR_MUTED, align=PP_ALIGN.LEFT)

    options = [
        {
            "title": "Project Online Sync",
            "badge": [("Allocations ONLY", CLR_ORANGE)],
            "pros": [
                "Automatisierte, geplante Synchronisation",
                "Admin-gesteuerte Governance",
                "Zentralisierter Prozess",
            ],
            "cons": [
                "Keine Schedule-Daten importierbar",
                "Root-URL & globale API-Credentials nötig",
            ],
            "accent": CLR_ORANGE,
            "icon": "⟳",
        },
        {
            "title": "MSP Add-In",
            "badge": [("Allocations ONLY", CLR_ORANGE)],
            "pros": [
                "Volle Kontrolle beim Project Manager",
                "Datenvalidierung & Vorschau vor Sync",
                "Erweitertes WBS-Mapping",
            ],
            "cons": [
                "Manuelle Ausführung pro Projekt",
                "Individuelle API-Token je Nutzer",
            ],
            "accent": CLR_ORANGE,
            "icon": "⎇",
        },
        {
            "title": "MPP Import",
            "badge": [("Full Schedule", CLR_BLUE), ("Allocations", CLR_ORANGE)],
            "pros": [
                "Einziges Tool für vollständige Schedules\n(Tasks, Dependencies, Dauer, Termine)",
                "Ideal für initiale Systemmigrationen",
            ],
            "cons": [
                "Manueller, dateibasierter Drag-&-Drop-Prozess",
                "Überschreibt bestehende Schedule-Daten",
            ],
            "accent": CLR_BLUE,
            "icon": "⇪",
        },
    ]

    card_y = 1.55
    card_w = 4.05
    card_h = 5.60
    gap = 0.12
    start_x = 0.5

    for i, opt in enumerate(options):
        x = start_x + i * (card_w + gap)

        # Karte
        card = add_rect(slide, x, card_y, card_w, card_h,
                        fill=CLR_CARD, line=CLR_BORDER)

        # Akzent-Kopfbalken
        head = add_rect(slide, x, card_y, card_w, 0.95,
                        fill=opt["accent"], line=None)
        head.adjustments[0] = 0.06
        set_text(head, f'{opt["icon"]}   {opt["title"]}',
                 size=16, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))

        # Badges
        bx = x + 0.35
        by = card_y + 1.10
        for j, (label, col) in enumerate(opt["badge"]):
            # Breite grob auf Textlänge
            bw = 1.55 if len(label) < 15 else 2.0
            add_pill(slide, bx, by, bw, 0.40, label,
                     fill=col, size=10.5)
            bx += bw + 0.15

        # Pros
        py = card_y + 1.80
        add_text(slide, x + 0.25, py, card_w - 0.5, 0.35,
                 "Vorteile", size=11.5, bold=True, color=CLR_GREEN,
                 align=PP_ALIGN.LEFT)
        tb = slide.shapes.add_textbox(Inches(x + 0.25), Inches(py + 0.35),
                                      Inches(card_w - 0.5), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        for k, line in enumerate(opt["pros"]):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = "✓ " + line
            run.font.size = Pt(10.5)
            run.font.color.rgb = CLR_TEXT
            run.font.name = "Calibri"
            p.space_after = Pt(4)

        # Cons
        cyy = card_y + 4.00
        add_text(slide, x + 0.25, cyy, card_w - 0.5, 0.35,
                 "Nachteile", size=11.5, bold=True, color=CLR_RED,
                 align=PP_ALIGN.LEFT)
        tb = slide.shapes.add_textbox(Inches(x + 0.25), Inches(cyy + 0.35),
                                      Inches(card_w - 0.5), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        for k, line in enumerate(opt["cons"]):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = "! " + line
            run.font.size = Pt(10.5)
            run.font.color.rgb = CLR_TEXT
            run.font.name = "Calibri"
            p.space_after = Pt(4)


# ---------------------------------------------------------------------------
# SLIDE 3 – Selecting Your Integration Strategy
# ---------------------------------------------------------------------------
def build_slide_3a(prs):
    """Variante A – 3 Karten mit Icon-Kreis, Situation, Pfeil, Tool (Nachbau)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CLR_BG)

    slide_title_runs(slide, [
        ("Selecting Your Integration Strategy", True, CLR_DARK),
    ], y=0.35, size=28)

    cards = [
        {
            "icon": "⇶",
            "situation": "If your operational model is…",
            "body": "Centralised planning requiring automated, hands-off "
                    "synchronisation across the enterprise.",
            "arrow_text": "Then deploy",
            "tool_l1": "Project",
            "tool_l2": "Online Sync",
        },
        {
            "icon": "👥",
            "situation": "If your operational model is…",
            "body": "Decentralised planning where individual Project Managers "
                    "must control, preview, and validate their own data "
                    "before submission.",
            "arrow_text": "Then deploy the",
            "tool_l1": "MSP",
            "tool_l2": "Add-In",
        },
        {
            "icon": "💾",
            "situation": "If your operational model requires…",
            "body": "A one-time system migration, or an ongoing necessity to "
                    "view full task logic, dependencies, and granular "
                    "project schedules in Tempus.",
            "arrow_text": "Then utilise",
            "tool_l1": "MPP",
            "tool_l2": "Import",
        },
    ]

    card_y = 1.25
    card_w = 4.05
    card_h = 5.80
    gap = 0.12
    start_x = 0.5

    for i, c in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rect(slide, x, card_y, card_w, card_h,
                 fill=CLR_CARD, line=CLR_BORDER)

        # Icon-Kreis
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(x + card_w / 2 - 0.75),
                                      Inches(card_y + 0.35),
                                      Inches(1.5), Inches(1.5))
        set_fill(circ, CLR_DARK)
        no_line(circ)
        set_text(circ, c["icon"], size=36, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))

        # Situation
        add_text(slide, x + 0.25, card_y + 2.10, card_w - 0.5, 0.4,
                 c["situation"], size=12, bold=True, color=CLR_DARK,
                 align=PP_ALIGN.CENTER)

        # Body
        add_text(slide, x + 0.25, card_y + 2.55, card_w - 0.5, 1.70,
                 c["body"], size=11, color=CLR_MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

        # Abwärtspfeil
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     Inches(x + card_w / 2 - 0.30),
                                     Inches(card_y + 4.15),
                                     Inches(0.60), Inches(0.55))
        set_fill(arr, CLR_DARK)
        no_line(arr)

        # Then deploy…
        add_text(slide, x + 0.25, card_y + 4.75, card_w - 0.5, 0.35,
                 c["arrow_text"], size=11, bold=True, color=CLR_DARK,
                 align=PP_ALIGN.CENTER)

        # Tool-Name groß
        add_text(slide, x + 0.25, card_y + 5.05, card_w - 0.5, 0.45,
                 c["tool_l1"], size=20, bold=True, color=CLR_DARK,
                 align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.25, card_y + 5.45, card_w - 0.5, 0.45,
                 c["tool_l2"], size=20, bold=True, color=CLR_DARK,
                 align=PP_ALIGN.CENTER)


def build_slide_3b(prs):
    """Variante B – Entscheidungsmatrix / Tabelle 'Wenn – Dann'."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CLR_CARD)

    slide_title_runs(slide, [
        ("Entscheidungs-Matrix: ", True, CLR_DARK),
        ("Welches Tool passt wann?", True, CLR_ORANGE),
    ], y=0.35, size=26)

    # Kopfzeile
    top = 1.35
    left = 0.5
    col_w = [4.5, 5.5, 2.7]
    row_h_head = 0.55
    row_heights = [1.45, 1.45, 1.45]

    def cell(x, y, w, h, fill, line=CLR_BORDER, text="", size=12, bold=False,
             color=CLR_TEXT, align=PP_ALIGN.LEFT):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
        set_fill(s, fill)
        set_line(s, line, 0.75)
        if text:
            set_text(s, text, size=size, bold=bold, color=color, align=align)
        return s

    # Header
    cell(left,              top, col_w[0], row_h_head,
         fill=CLR_DARK, line=CLR_DARK,
         text="Operativer Rahmen", size=13, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.LEFT)
    cell(left + col_w[0],   top, col_w[1], row_h_head,
         fill=CLR_DARK, line=CLR_DARK,
         text="Merkmale", size=13, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.LEFT)
    cell(left + col_w[0] + col_w[1], top, col_w[2], row_h_head,
         fill=CLR_DARK, line=CLR_DARK,
         text="Empfohlenes Tool", size=13, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    rows = [
        {
            "op": "Zentrale Planung",
            "desc": "Automatisierte, hands-off Synchronisation unternehmensweit. "
                    "Governance liegt beim Admin.",
            "tool": "Project\nOnline Sync",
            "accent": CLR_ORANGE,
        },
        {
            "op": "Dezentrale Planung",
            "desc": "Einzelne Project Manager steuern, prüfen und validieren "
                    "ihre Daten selbst vor dem Transfer.",
            "tool": "MSP\nAdd-In",
            "accent": CLR_ORANGE,
        },
        {
            "op": "Migration / Volle Tiefe",
            "desc": "Einmalige Systemmigration oder laufender Bedarf an vollständiger "
                    "Task-Logik, Dependencies und Terminen in Tempus.",
            "tool": "MPP\nImport",
            "accent": CLR_BLUE,
        },
    ]

    y = top + row_h_head
    for i, r in enumerate(rows):
        bg = CLR_BG if i % 2 == 0 else CLR_CARD
        h = row_heights[i]
        cell(left, y, col_w[0], h, fill=bg,
             text=r["op"], size=14, bold=True, color=CLR_DARK,
             align=PP_ALIGN.LEFT)
        cell(left + col_w[0], y, col_w[1], h, fill=bg,
             text=r["desc"], size=11.5, color=CLR_MUTED,
             align=PP_ALIGN.LEFT)
        # Tool-Zelle farbig
        cell(left + col_w[0] + col_w[1], y, col_w[2], h, fill=r["accent"],
             line=r["accent"],
             text=r["tool"], size=16, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)
        y += h

    # Footer
    add_text(slide, 0.5, 6.75, 12.3, 0.5,
             "Auswahlregel: Das Zielsystem (Allocations vs. Full Schedule) "
             "und das operative Modell bestimmen das Integrationswerkzeug.",
             size=11, bold=True, color=CLR_DARK, align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# TITEL- & AGENDA-FOLIE
# ---------------------------------------------------------------------------
def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CLR_DARK)

    add_text(slide, 0.8, 2.6, 12.0, 0.6,
             "TEMPUS × MS PROJECT", size=16, bold=True,
             color=CLR_ORANGE, align=PP_ALIGN.LEFT)
    add_text(slide, 0.8, 3.2, 12.0, 1.2,
             "Integrationsoptionen", size=48, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.LEFT)
    add_text(slide, 0.8, 4.35, 12.0, 0.8,
             "Editierbare Folienvarianten für den Unternehmens-Foliensatz",
             size=18, color=RGBColor(0xCF, 0xD4, 0xDB), align=PP_ALIGN.LEFT)

    # Akzent-Balken
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.8), Inches(3.10),
                                 Inches(0.12), Inches(1.10))
    set_fill(bar, CLR_ORANGE)
    no_line(bar)


def build_agenda_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CLR_CARD)

    slide_title_runs(slide, [("Inhalt", True, CLR_DARK)], y=0.4, size=30)

    items = [
        ("1", "Technische Abgrenzung", "Schedule vs. Allocation",
         "Slide 1A (Nachbau) · Slide 1B (Corporate)"),
        ("2", "Integrationsoptionen", "At a Glance",
         "Slide 2A (Tabelle) · Slide 2B (3-Karten)"),
        ("3", "Strategie-Auswahl", "Welches Tool wann?",
         "Slide 3A (Karten) · Slide 3B (Matrix)"),
    ]
    y = 1.5
    for num, title, sub, variants in items:
        # Nummer-Kreis
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(0.8), Inches(y),
                                      Inches(0.9), Inches(0.9))
        set_fill(circ, CLR_ORANGE)
        no_line(circ)
        set_text(circ, num, size=22, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
        # Titel
        add_text(slide, 2.0, y, 10.5, 0.45,
                 title + "  –  " + sub,
                 size=18, bold=True, color=CLR_DARK, align=PP_ALIGN.LEFT)
        add_text(slide, 2.0, y + 0.45, 10.5, 0.40,
                 variants, size=12, color=CLR_MUTED, align=PP_ALIGN.LEFT)
        y += 1.55


# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_title_slide(prs)
    build_agenda_slide(prs)

    build_slide_1a(prs)
    build_slide_1b(prs)

    build_slide_2a(prs)
    build_slide_2b(prs)

    build_slide_3a(prs)
    build_slide_3b(prs)

    out = "PowerPoint-Exports/Integration-Optionen_editierbar.pptx"
    prs.save(out)
    print(f"OK -> {out}")


if __name__ == "__main__":
    main()
