"""
Valkeen-Design · MODERN & EINHEITLICH
=====================================

Ziel: Ein durchgängiges, minimalistisches Folien-Set mit nur leichten
Farbakzenten.

Design-Regeln (bewusst streng):
  - Hintergrund:            immer reinweiß  (#FFFFFF)
  - Primärfarbe:            Valkeen-Navy   (#002856)  – nur für Titel,
                            Quellblock, kleine Typo-Akzente
  - Akzent 1:               Valkeen-Green  (#86BC25)  – ausschließlich als
                            3-px-Akzentstrich (Titel-Balken, Card-Top-Border)
  - Akzent 2 (Pastell):     Soft-Orange    (#FDE9D3)  – Allocations-Fläche
                            Soft-Slate     (#E4ECF3)  – Full-Schedule-Fläche
  - Linien/Rahmen:          #E5E7EB  (1 px)
  - Text dunkel:            #1C2835
  - Text mittel/muted:      #6B7280
  - Typografie:             Calibri Light (Titel, Tool-Names)
                            Calibri (Body)
  - KEINE kräftigen Farbflächen, KEINE Schatten, KEINE gerundeten Pills
    (stattdessen rechteckige Chips mit dünner Outline).
  - Jede Folie nutzt dasselbe Title-Grid:
        Akzentstrich bei (x=0.5, y=0.50) · Titel bei (x=0.70, y=0.40)
        Divider-Linie bei y=1.15

Ausgaben:
  PowerPoint-Exports/Valkeen_Integration-Optionen_MODERN.pptx
  PowerPoint-Exports/Entscheidungsvorlage_Knauf_mit_Integrationsfolien_MODERN.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Reduzierte Palette
# ---------------------------------------------------------------------------
NAVY       = RGBColor(0x00, 0x28, 0x56)
GREEN      = RGBColor(0x86, 0xBC, 0x25)
SOFT_ORG   = RGBColor(0xFD, 0xE9, 0xD3)  # pastell
SOFT_SLATE = RGBColor(0xE4, 0xEC, 0xF3)  # pastell
BORDER     = RGBColor(0xE5, 0xE7, 0xEB)
BORDER_DK  = RGBColor(0xCB, 0xD0, 0xD8)
TEXT_DARK  = RGBColor(0x1C, 0x28, 0x35)
TEXT_MED   = RGBColor(0x4B, 0x55, 0x63)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Calibri Light"
BODY_FONT  = "Calibri"

SLIDE_W, SLIDE_H = 13.333, 7.5


# ---------------------------------------------------------------------------
# Tiny helpers
# ---------------------------------------------------------------------------
def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_fill(shape):
    shape.fill.background()


def _line(shape, color, width_pt=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def _no_line(shape):
    shape.line.fill.background()


def _kill_shadow(shape):
    spPr = shape._element.spPr
    for tag in ('a:effectLst', 'a:effectRef'):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))


def _text(shape, text, *, size=12, bold=False, color=TEXT_DARK,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
          font=BODY_FONT, italic=False, line_spacing=None,
          letter_spacing=None):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    tf.clear()
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        if letter_spacing:
            # Character spacing in 1/100 pt
            rPr = run._r.get_or_add_rPr()
            rPr.set('spc', str(letter_spacing))


def add_rect(slide, x, y, w, h, *, fill=WHITE, line=BORDER, line_w=0.75,
             rounded=False, radius=0.03):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    if rounded:
        s.adjustments[0] = radius
    if fill is None:
        _no_fill(s)
    else:
        _fill(s, fill)
    if line is None:
        _no_line(s)
    else:
        _line(s, line, line_w)
    _kill_shadow(s)
    return s


def add_text(slide, x, y, w, h, text, *, size=12, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=BODY_FONT,
             italic=False, line_spacing=None, letter_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _text(tb, text, size=size, bold=bold, color=color, align=align,
          anchor=anchor, font=font, italic=italic,
          line_spacing=line_spacing, letter_spacing=letter_spacing)
    return tb


def add_hline(slide, x1, y, x2, *, color=BORDER, width_pt=0.75):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y),
                                    Inches(x2), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(width_pt)
    return ln


def add_vline(slide, x, y1, y2, *, color=BORDER, width_pt=0.75):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x), Inches(y1),
                                    Inches(x), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width_pt)
    return ln


def add_chip(slide, x, y, w, h, text, *, fill=None, line=BORDER_DK,
             text_color=TEXT_DARK, size=10, bold=True, letter_spacing=30):
    """Rechteckiger Chip mit dünner Outline (kein Pill-Look)."""
    s = add_rect(slide, x, y, w, h, fill=fill if fill else WHITE,
                 line=line, line_w=0.75, rounded=False)
    _text(s, text, size=size, bold=bold, color=text_color, font=BODY_FONT,
          letter_spacing=letter_spacing)
    return s


# ---------------------------------------------------------------------------
# Einheitliche Title-Bar
# ---------------------------------------------------------------------------
def add_title(slide, title, subtitle=None, *, eyebrow=None):
    # Eyebrow (Kleingedrucktes über dem Titel, z. B. Kategorie)
    if eyebrow:
        add_text(slide, 0.70, 0.32, 12.0, 0.25, eyebrow.upper(),
                 size=9, bold=True, color=GREEN, font=BODY_FONT,
                 letter_spacing=200, align=PP_ALIGN.LEFT)
        title_y = 0.55
    else:
        title_y = 0.42

    # Grüner Akzentstrich links
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.50), Inches(title_y + 0.10),
                                 Inches(0.04), Inches(0.50))
    _fill(bar, GREEN)
    _no_line(bar)
    _kill_shadow(bar)

    # Titel
    add_text(slide, 0.70, title_y, 12.2, 0.60, title,
             size=28, bold=False, color=NAVY, font=TITLE_FONT,
             align=PP_ALIGN.LEFT, line_spacing=1.0)

    # Subtitle
    if subtitle:
        add_text(slide, 0.70, title_y + 0.60, 12.2, 0.35, subtitle,
                 size=12, color=TEXT_MUTED, font=BODY_FONT,
                 align=PP_ALIGN.LEFT, italic=False)

    # Divider – unter dem Titel
    add_hline(slide, 0.50, 1.30, 12.83, color=BORDER, width_pt=0.5)


def add_footer(slide, page_num, total):
    add_hline(slide, 0.50, 7.08, 12.83, color=BORDER, width_pt=0.5)
    add_text(slide, 0.50, 7.15, 10.0, 0.26,
             "Valkeen GmbH  ·  MS Project ⇄ Tempus Integration",
             size=8.5, color=TEXT_MUTED, align=PP_ALIGN.LEFT,
             font=BODY_FONT, letter_spacing=30)
    add_text(slide, 11.30, 7.15, 1.50, 0.26,
             f"{page_num:02d} / {total:02d}",
             size=8.5, color=NAVY, align=PP_ALIGN.RIGHT, bold=True,
             font=BODY_FONT, letter_spacing=30)


def set_slide_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, WHITE)
    return s


# ===========================================================================
# SLIDE – Title-Cover
# ===========================================================================
def slide_cover(prs):
    slide = new_slide(prs)
    # sehr zurückhaltendes Cover: links Titel, rechts dezente Grafik
    # Eyebrow
    add_text(slide, 0.80, 2.85, 11.5, 0.30,
             "TEMPUS · MS PROJECT INTEGRATION",
             size=10, bold=True, color=GREEN, align=PP_ALIGN.LEFT,
             font=BODY_FONT, letter_spacing=300)
    # grüner Akzentstrich
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.80), Inches(3.25),
                                 Inches(0.70), Inches(0.04))
    _fill(bar, GREEN)
    _no_line(bar)
    _kill_shadow(bar)
    # Titel
    add_text(slide, 0.80, 3.40, 11.5, 1.20, "Integrationsoptionen",
             size=52, bold=False, color=NAVY, font=TITLE_FONT,
             align=PP_ALIGN.LEFT, line_spacing=1.0)
    add_text(slide, 0.80, 4.40, 11.5, 0.55,
             "Schedule vs. Allocation — ein einheitlicher Überblick.",
             size=16, color=TEXT_MUTED, font=BODY_FONT, align=PP_ALIGN.LEFT)
    # Footer
    add_hline(slide, 0.80, 6.80, 12.53, color=BORDER)
    add_text(slide, 0.80, 6.88, 6.0, 0.30, "Valkeen GmbH",
             size=10, bold=True, color=NAVY, font=BODY_FONT,
             letter_spacing=80)
    add_text(slide, 6.53, 6.88, 6.0, 0.30, "www.valkeen.com",
             size=10, color=TEXT_MUTED, font=BODY_FONT,
             align=PP_ALIGN.RIGHT, letter_spacing=80)


# ===========================================================================
# SLIDE – Agenda
# ===========================================================================
def slide_agenda(prs, page_num, total):
    slide = new_slide(prs)
    add_title(slide, "Inhalt", eyebrow="Agenda")

    items = [
        ("01", "Technical Boundary",
         "Schedule vs. Allocation — das Zielsystem bestimmt das Werkzeug."),
        ("02", "Integration Options",
         "Drei Werkzeuge im direkten Vergleich: Stärken & Limitierungen."),
        ("03", "Selecting Your Strategy",
         "Operatives Modell → passendes Integrationswerkzeug."),
    ]
    y = 1.90
    for num, title, sub in items:
        add_text(slide, 0.70, y, 0.90, 0.50, num,
                 size=22, bold=False, color=NAVY, font=TITLE_FONT,
                 align=PP_ALIGN.LEFT)
        add_vline(slide, 1.55, y + 0.05, y + 0.85, color=BORDER)
        add_text(slide, 1.75, y + 0.00, 11.0, 0.45, title,
                 size=18, bold=False, color=NAVY, font=TITLE_FONT)
        add_text(slide, 1.75, y + 0.50, 11.0, 0.35, sub,
                 size=11, color=TEXT_MUTED, font=BODY_FONT)
        y += 1.20

    add_footer(slide, page_num, total)


# ===========================================================================
# SLIDE 1 — Technical Boundary (modern, einheitlich)
# ===========================================================================
def slide_boundary(prs, page_num, total):
    slide = new_slide(prs)
    add_title(slide,
              "The Technical Boundary",
              subtitle="Das Zielsystem der Daten bestimmt das Integrationswerkzeug.",
              eyebrow="01 · Schedule vs. Allocation")

    # --- Source --------------------------------------------------------
    src_x, src_y, src_w, src_h = 0.50, 3.60, 2.00, 1.10
    src = add_rect(slide, src_x, src_y, src_w, src_h,
                   fill=WHITE, line=BORDER_DK, line_w=0.75)
    # Dünner Navy-Top-Border als einziger Farbakzent
    top_bar = add_rect(slide, src_x, src_y, src_w, 0.05,
                       fill=NAVY, line=None)
    _text(src, "MS Project\nData", size=14, bold=False, color=NAVY,
          font=TITLE_FONT)

    # --- Geometrie ------------------------------------------------------
    top_y = 2.30
    bot_y = 4.90
    lane_left = 2.95
    lane_right = 10.60  # left border of target block
    lane_h = 1.55

    # --- LANE 1 · Allocations (pastell-orange) --------------------------
    add_rect(slide, lane_left, top_y, lane_right - lane_left, lane_h,
             fill=SOFT_ORG, line=None)
    # Top accent line
    add_hline(slide, lane_left, top_y, lane_right, color=GREEN, width_pt=1.0)
    # Label links
    add_text(slide, lane_left + 0.20, top_y + 0.10, 2.5, 0.30,
             "PATH A",
             size=9, bold=True, color=NAVY, font=BODY_FONT,
             letter_spacing=200)
    add_text(slide, lane_left + 0.20, top_y + 0.35, 3.5, 0.30,
             "Allocation-only tools",
             size=11, bold=True, color=TEXT_DARK, font=BODY_FONT)

    # Chips – 3 Stück, rechts in der Lane
    chip_y = top_y + 0.85
    chip_h = 0.50
    chip_w = 1.75
    chip_gap = 0.12
    chips_total_w = 3 * chip_w + 2 * chip_gap
    chip_start = lane_left + (lane_right - lane_left) - chips_total_w - 0.20
    for i, label in enumerate(["Project Online Sync", "MSP Add-In", "MPP Import"]):
        add_chip(slide, chip_start + i * (chip_w + chip_gap), chip_y,
                 chip_w, chip_h, label, fill=WHITE, line=BORDER_DK,
                 text_color=NAVY, size=10, bold=True, letter_spacing=30)

    # Ziel: Allocations Grid
    tgt_x = lane_right + 0.20
    add_rect(slide, tgt_x, top_y, 2.10, lane_h,
             fill=WHITE, line=BORDER_DK)
    add_hline(slide, tgt_x, top_y, tgt_x + 2.10, color=GREEN, width_pt=1.0)
    add_text(slide, tgt_x + 0.15, top_y + 0.20, 1.80, 0.35,
             "TARGET", size=9, bold=True, color=NAVY,
             letter_spacing=200, font=BODY_FONT)
    add_text(slide, tgt_x + 0.15, top_y + 0.45, 1.80, 0.45,
             "Allocations Grid",
             size=14, bold=False, color=NAVY, font=TITLE_FONT)
    add_text(slide, tgt_x + 0.15, top_y + 0.95, 1.80, 0.55,
             "High-level planning & resource demand.",
             size=9.5, color=TEXT_MUTED, font=BODY_FONT,
             line_spacing=1.15)

    # --- LANE 2 · Full Schedule (pastell-slate) -------------------------
    add_rect(slide, lane_left, bot_y, lane_right - lane_left, lane_h,
             fill=SOFT_SLATE, line=None)
    add_hline(slide, lane_left, bot_y, lane_right, color=NAVY, width_pt=1.0)

    add_text(slide, lane_left + 0.20, bot_y + 0.10, 2.5, 0.30,
             "PATH B",
             size=9, bold=True, color=NAVY, font=BODY_FONT,
             letter_spacing=200)
    add_text(slide, lane_left + 0.20, bot_y + 0.35, 3.5, 0.30,
             "Full-schedule tool",
             size=11, bold=True, color=TEXT_DARK, font=BODY_FONT)

    # Ein Chip (MPP Import) – mittig
    chip_w2 = 1.75
    cx = lane_left + (lane_right - lane_left) - chips_total_w - 0.20 \
         + (chip_w + chip_gap)  # mittlere Position wie oben
    add_chip(slide, cx, bot_y + 0.85, chip_w2, chip_h, "MPP Import",
             fill=WHITE, line=BORDER_DK, text_color=NAVY, size=10, bold=True,
             letter_spacing=30)

    # Ziel: Full Schedule
    add_rect(slide, tgt_x, bot_y, 2.10, lane_h, fill=WHITE, line=BORDER_DK)
    add_hline(slide, tgt_x, bot_y, tgt_x + 2.10, color=NAVY, width_pt=1.0)
    add_text(slide, tgt_x + 0.15, bot_y + 0.20, 1.80, 0.35,
             "TARGET", size=9, bold=True, color=NAVY,
             letter_spacing=200, font=BODY_FONT)
    add_text(slide, tgt_x + 0.15, bot_y + 0.45, 1.80, 0.45,
             "Full Schedule",
             size=14, bold=False, color=NAVY, font=TITLE_FONT)
    add_text(slide, tgt_x + 0.15, bot_y + 0.95, 1.80, 0.55,
             "Tasks, dependencies, dates & durations.",
             size=9.5, color=TEXT_MUTED, font=BODY_FONT,
             line_spacing=1.15)

    # Verbindungs-Linien (schlank, einheitlich grau)
    # Source (rechter Rand) -> Lane A (linker Rand, vertikal mittig)
    x_src_right = src_x + src_w
    y_src_mid = src_y + src_h / 2
    # Lane A
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x_src_right), Inches(y_src_mid),
                                    Inches(lane_left), Inches(top_y + lane_h / 2))
    ln.line.color.rgb = BORDER_DK
    ln.line.width = Pt(1.25)
    # Pfeilspitze
    lnXml = ln.line._get_or_add_ln()
    tail = etree.SubElement(lnXml, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'sm'); tail.set('h', 'sm')

    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x_src_right), Inches(y_src_mid),
                                    Inches(lane_left), Inches(bot_y + lane_h / 2))
    ln.line.color.rgb = BORDER_DK
    ln.line.width = Pt(1.25)
    lnXml = ln.line._get_or_add_ln()
    tail = etree.SubElement(lnXml, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'sm'); tail.set('h', 'sm')

    # --- Golden Rule Callout (ganz unten rechts, sehr dezent) ----------
    add_text(slide, 0.50, 6.75, 12.3, 0.28,
             "Merksatz · Das Zielsystem der Daten entscheidet, welches "
             "Integrationswerkzeug zum Einsatz kommt.",
             size=10, italic=True, color=TEXT_MUTED, font=BODY_FONT,
             align=PP_ALIGN.LEFT)

    add_footer(slide, page_num, total)


# ===========================================================================
# SLIDE 2 — Integration Options (modern Table, nur horizontale Divider)
# ===========================================================================
def slide_options(prs, page_num, total):
    slide = new_slide(prs)
    add_title(slide,
              "Integration Options",
              subtitle="Drei Werkzeuge — Zielsystem, Stärken und Limitierungen im Überblick.",
              eyebrow="02 · At a Glance")

    # --- Table Geometry -------------------------------------------------
    tbl_left = 0.50
    tbl_top  = 1.65
    lbl_w    = 1.30
    col_w    = (13.33 - 1.00 - lbl_w) / 3    # exakt gleich breit
    total_w  = lbl_w + 3 * col_w

    # Headers (Tool-Namen)
    headers = ["Project Online Sync", "MSP Add-In", "MPP Import"]
    head_h  = 0.70

    for i, h in enumerate(headers):
        x = tbl_left + lbl_w + i * col_w
        add_text(slide, x + 0.10, tbl_top + 0.15, col_w - 0.20, 0.40, h,
                 size=15, bold=False, color=NAVY, font=TITLE_FONT,
                 align=PP_ALIGN.LEFT)

    # horizontale Divider unter Headern (1 px Navy) – einziger farbiger Strich
    add_hline(slide, tbl_left, tbl_top + head_h, tbl_left + total_w,
              color=NAVY, width_pt=1.25)

    # vertikale dünne Divider zwischen den Spalten
    for i in range(1, 4):
        add_vline(slide, tbl_left + lbl_w + (i - 1) * col_w if i == 1 else tbl_left + lbl_w + (i - 1) * col_w,
                  tbl_top + 0.10, tbl_top + 5.10, color=BORDER)

    # === Row 1: Data Target ============================================
    row1_y = tbl_top + head_h + 0.15
    row1_h = 0.80
    add_text(slide, tbl_left + 0.05, row1_y + 0.10, lbl_w - 0.10, 0.30,
             "DATA TARGET", size=9, bold=True, color=TEXT_MUTED,
             letter_spacing=200, font=BODY_FONT)

    # Chip-Definitionen je Spalte
    targets = [
        [("Allocations only", SOFT_ORG)],
        [("Allocations only", SOFT_ORG)],
        [("Full Schedule", SOFT_SLATE), ("Allocations", SOFT_ORG)],
    ]
    for i, chips in enumerate(targets):
        x = tbl_left + lbl_w + i * col_w
        cx = x + 0.15
        cy = row1_y + 0.22
        for label, fill in chips:
            # Textbreite schätzen: ~ 0.10 in pro Zeichen bei 10pt bold
            bw = 0.18 + 0.08 * len(label)
            bw = max(1.55, min(bw, col_w - 0.30))
            add_chip(slide, cx, cy, bw, 0.38, label,
                     fill=fill, line=None,
                     text_color=NAVY, size=10, bold=True, letter_spacing=30)
            cx += bw + 0.12

    add_hline(slide, tbl_left, row1_y + row1_h + 0.15,
              tbl_left + total_w, color=BORDER)

    # === Row 2: STRENGTHS ==============================================
    row2_y = row1_y + row1_h + 0.25
    row2_h = 1.90
    add_text(slide, tbl_left + 0.05, row2_y + 0.05, lbl_w - 0.10, 0.30,
             "STRENGTHS", size=9, bold=True, color=TEXT_MUTED,
             letter_spacing=200, font=BODY_FONT)

    pros = [
        [
            "Automatisierte, geplante Synchronisation",
            "Admin-gesteuerte Governance",
            "Zentralisierter Prozess",
        ],
        [
            "Volle Kontrolle beim Project Manager",
            "Datenvalidierung & Fehlervorschau vor Sync",
            "Erweitertes WBS-Mapping",
        ],
        [
            "Einziges Tool für vollständige Schedules\n(Tasks, Dependencies, Dauer, Termine)",
            "Ideal für initiale Systemmigrationen",
        ],
    ]
    for i, items in enumerate(pros):
        x = tbl_left + lbl_w + i * col_w
        tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(row2_y + 0.05),
                                      Inches(col_w - 0.30), Inches(row2_h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        for k, line in enumerate(items):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)
            # dezenter grüner Marker (Punkt) statt Häkchen
            r1 = p.add_run()
            r1.text = "•  "
            r1.font.size = Pt(13)
            r1.font.bold = True
            r1.font.color.rgb = GREEN
            r1.font.name = BODY_FONT
            r2 = p.add_run()
            r2.text = line
            r2.font.size = Pt(10.5)
            r2.font.color.rgb = TEXT_DARK
            r2.font.name = BODY_FONT

    add_hline(slide, tbl_left, row2_y + row2_h + 0.10,
              tbl_left + total_w, color=BORDER)

    # === Row 3: CONSIDERATIONS =========================================
    row3_y = row2_y + row2_h + 0.20
    row3_h = 1.60
    add_text(slide, tbl_left + 0.05, row3_y + 0.05, lbl_w - 0.10, 0.30,
             "CONSIDERATIONS", size=9, bold=True, color=TEXT_MUTED,
             letter_spacing=200, font=BODY_FONT)

    cons = [
        [
            "Keine Schedule-Daten importierbar",
            "Root-URL & globale API-Credentials nötig",
        ],
        [
            "Manuelle Ausführung pro Projekt",
            "Individuelle API-Token je Nutzer",
        ],
        [
            "Manueller Drag-&-Drop-Prozess (dateibasiert)",
            "Überschreibt bestehende Schedule-Daten",
        ],
    ]
    for i, items in enumerate(cons):
        x = tbl_left + lbl_w + i * col_w
        tb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(row3_y + 0.05),
                                      Inches(col_w - 0.30), Inches(row3_h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        for k, line in enumerate(items):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)
            r1 = p.add_run()
            r1.text = "—  "
            r1.font.size = Pt(11)
            r1.font.bold = True
            r1.font.color.rgb = NAVY
            r1.font.name = BODY_FONT
            r2 = p.add_run()
            r2.text = line
            r2.font.size = Pt(10.5)
            r2.font.color.rgb = TEXT_MED
            r2.font.name = BODY_FONT

    add_footer(slide, page_num, total)


# ===========================================================================
# SLIDE 3 — Selecting Your Strategy (3 klare Karten, minimale Akzente)
# ===========================================================================
def _mini_icon(slide, kind, cx, cy, size=0.42):
    """Dezentes Mono-Icon in Navy – rein aus einfachen Shapes."""
    # Hintergrundkreis / Akzentkreis (dünne Outline, kein Fill)
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(cx - size), Inches(cy - size),
                                  Inches(size * 2), Inches(size * 2))
    _fill(ring, WHITE)
    _line(ring, BORDER_DK, 0.75)
    _kill_shadow(ring)

    if kind == "sync":
        # kleiner gebogener Pfeil (Kreispfeil) als Zeichen für "Sync"
        arr = slide.shapes.add_shape(MSO_SHAPE.CIRCULAR_ARROW,
                                     Inches(cx - size * 0.55),
                                     Inches(cy - size * 0.55),
                                     Inches(size * 1.1),
                                     Inches(size * 1.1))
        _fill(arr, NAVY)
        _no_line(arr)
        _kill_shadow(arr)
    elif kind == "hub":
        # zentraler Punkt + 5 kleine Punkte drumherum
        import math
        dot_r = size * 0.20
        # Zentrum
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(cx - dot_r), Inches(cy - dot_r),
                                   Inches(dot_r * 2), Inches(dot_r * 2))
        _fill(d, NAVY); _no_line(d); _kill_shadow(d)
        # 5 kleine Punkte
        R = size * 0.70
        small = size * 0.13
        for k in range(5):
            ang = math.radians(72 * k - 90)
            sx = cx + R * math.cos(ang)
            sy = cy + R * math.sin(ang)
            # Verbindungslinie ganz dünn
            ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                            Inches(cx), Inches(cy),
                                            Inches(sx), Inches(sy))
            ln.line.color.rgb = BORDER_DK
            ln.line.width = Pt(0.5)
            dn = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(sx - small),
                                        Inches(sy - small),
                                        Inches(small * 2), Inches(small * 2))
            _fill(dn, NAVY); _no_line(dn); _kill_shadow(dn)
    elif kind == "disk":
        # 3 übereinanderliegende schmale Ovale (Disk-Stack)
        w = size * 1.20
        h_disk = size * 0.28
        spacing = size * 0.30
        for i in range(3):
            dy = cy - spacing + i * spacing
            d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(cx - w / 2), Inches(dy - h_disk / 2),
                                       Inches(w), Inches(h_disk))
            _fill(d, NAVY); _no_line(d); _kill_shadow(d)


def slide_strategy(prs, page_num, total):
    slide = new_slide(prs)
    add_title(slide,
              "Selecting Your Strategy",
              subtitle="Das operative Modell entscheidet, welches Werkzeug passt.",
              eyebrow="03 · Decision Framework")

    cards = [
        {
            "icon": "sync",
            "ctx": "Centralised planning",
            "body": "Automatisierte, unternehmensweite Synchronisation — "
                    "hands-off, Admin-gesteuert.",
            "tool": "Project Online Sync",
            "accent": GREEN,
        },
        {
            "icon": "hub",
            "ctx": "Decentralised planning",
            "body": "Project Manager steuern, prüfen und validieren ihre "
                    "Daten selbst vor dem Transfer.",
            "tool": "MSP Add-In",
            "accent": GREEN,
        },
        {
            "icon": "disk",
            "ctx": "Migration · Full depth",
            "body": "Einmalige Migration oder laufender Bedarf an vollständiger "
                    "Task-Logik, Dependencies und Terminen.",
            "tool": "MPP Import",
            "accent": NAVY,
        },
    ]

    card_y = 1.70
    card_w = 4.00
    card_h = 5.10
    gap    = 0.22
    start_x = 0.50 + (13.333 - 1.0 - 3 * card_w - 2 * gap) / 2  # zentrieren

    for i, c in enumerate(cards):
        x = start_x + i * (card_w + gap)

        # Karte – weiß, dünne Outline
        add_rect(slide, x, card_y, card_w, card_h,
                 fill=WHITE, line=BORDER_DK, line_w=0.75)
        # Top-Akzent 3-px-Linie
        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(x), Inches(card_y),
                                     Inches(card_w), Inches(0.05))
        _fill(top, c["accent"])
        _no_line(top); _kill_shadow(top)

        # Step number
        add_text(slide, x + 0.35, card_y + 0.25, 1.0, 0.30,
                 f"STEP {i+1:02d}",
                 size=9, bold=True, color=c["accent"], font=BODY_FONT,
                 letter_spacing=250, align=PP_ALIGN.LEFT)

        # Icon
        _mini_icon(slide, c["icon"],
                   cx=x + card_w / 2, cy=card_y + 1.45, size=0.55)

        # Context-Label
        add_text(slide, x + 0.35, card_y + 2.30, card_w - 0.70, 0.40,
                 c["ctx"], size=16, bold=False, color=NAVY,
                 font=TITLE_FONT, align=PP_ALIGN.CENTER)

        # Body
        add_text(slide, x + 0.35, card_y + 2.80, card_w - 0.70, 1.35,
                 c["body"], size=11, color=TEXT_MED, font=BODY_FONT,
                 align=PP_ALIGN.CENTER, line_spacing=1.3)

        # Divider
        add_hline(slide, x + 0.65, card_y + 4.25,
                  x + card_w - 0.65, color=BORDER)

        # Deploy Label
        add_text(slide, x + 0.35, card_y + 4.38, card_w - 0.70, 0.25,
                 "DEPLOY", size=8.5, bold=True, color=TEXT_MUTED,
                 letter_spacing=250, font=BODY_FONT, align=PP_ALIGN.CENTER)

        # Tool Name (gross)
        add_text(slide, x + 0.25, card_y + 4.62, card_w - 0.50, 0.50,
                 c["tool"], size=18, bold=False, color=NAVY,
                 font=TITLE_FONT, align=PP_ALIGN.CENTER)

    add_footer(slide, page_num, total)


# ===========================================================================
# SLIDE 4 — Matrix-Alternative (optional, im gleichen modernen Stil)
# ===========================================================================
def slide_matrix(prs, page_num, total):
    slide = new_slide(prs)
    add_title(slide,
              "Decision Matrix",
              subtitle="Operativer Rahmen × Zielsystem → Empfehlung.",
              eyebrow="Alternative · Tabellarisch")

    # Geometrie
    left = 0.50
    top  = 1.65
    cols_w = [4.00, 5.30, 2.00, 1.03]  # Rahmen, Merkmale, Zielsystem, Tool
    header_h = 0.55

    # Header-Zeile (nur Text + Navy-Divider unten)
    xs = [left]
    for w in cols_w:
        xs.append(xs[-1] + w)
    headers = ["Operativer Rahmen", "Merkmale", "Zielsystem", "Tool"]
    for i, h in enumerate(headers):
        align = PP_ALIGN.LEFT if i < 2 else PP_ALIGN.CENTER
        add_text(slide, xs[i] + 0.08, top + 0.12, cols_w[i] - 0.16, 0.35,
                 h, size=10, bold=True, color=TEXT_MUTED,
                 letter_spacing=200, font=BODY_FONT, align=align)
    add_hline(slide, left, top + header_h, xs[-1], color=NAVY, width_pt=1.25)

    rows = [
        {
            "op": "Zentrale Planung",
            "desc": "Automatisierte, hands-off Synchronisation unternehmensweit. "
                    "Governance liegt beim Admin.",
            "target": "Allocations",
            "target_fill": SOFT_ORG,
            "tool": "Project Online Sync",
        },
        {
            "op": "Dezentrale Planung",
            "desc": "Einzelne Project Manager steuern, prüfen und validieren "
                    "ihre Daten vor dem Transfer.",
            "target": "Allocations",
            "target_fill": SOFT_ORG,
            "tool": "MSP Add-In",
        },
        {
            "op": "Migration · Full depth",
            "desc": "Einmalige Systemmigration oder laufender Bedarf an "
                    "vollständiger Task-Logik, Dependencies und Terminen.",
            "target": "Full Schedule",
            "target_fill": SOFT_SLATE,
            "tool": "MPP Import",
        },
    ]

    row_h = 1.45
    y = top + header_h + 0.05
    for i, r in enumerate(rows):
        # Op
        add_text(slide, xs[0] + 0.08, y + 0.30, cols_w[0] - 0.16, 0.50,
                 r["op"], size=14, bold=False, color=NAVY,
                 font=TITLE_FONT, align=PP_ALIGN.LEFT)
        # Desc
        add_text(slide, xs[1] + 0.08, y + 0.20, cols_w[1] - 0.16, row_h - 0.30,
                 r["desc"], size=10.5, color=TEXT_MED, font=BODY_FONT,
                 align=PP_ALIGN.LEFT, line_spacing=1.25)
        # Target-Chip (pastell)
        chip_w = 1.70
        chip_x = xs[2] + (cols_w[2] - chip_w) / 2
        add_chip(slide, chip_x, y + 0.48, chip_w, 0.42, r["target"],
                 fill=r["target_fill"], line=None, text_color=NAVY,
                 size=10.5, bold=True, letter_spacing=30)
        # Tool (rechts, kurz)
        add_text(slide, xs[3] + 0.04, y + 0.30, cols_w[3] - 0.08, 0.70,
                 r["tool"], size=11, bold=False, color=NAVY,
                 font=TITLE_FONT, align=PP_ALIGN.CENTER)
        # Divider
        add_hline(slide, left, y + row_h, xs[-1], color=BORDER)
        y += row_h

    add_footer(slide, page_num, total)


# ===========================================================================
# Build
# ===========================================================================
def build_deck(prs):
    total = 6  # Cover + Agenda + 3 Haupt + 1 Matrix
    slide_cover(prs)
    slide_agenda(prs, 2, total)
    slide_boundary(prs, 3, total)
    slide_options(prs, 4, total)
    slide_strategy(prs, 5, total)
    slide_matrix(prs, 6, total)


def build_standalone(out_path):
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    build_deck(prs)
    prs.save(out_path)
    print(f"OK · standalone -> {out_path}")


def build_on_template(template_path, out_path):
    prs = Presentation(template_path)
    start_idx = len(prs.slides) + 1
    total = start_idx + 3   # nur die 4 Content-Folien anhängen
    slide_boundary(prs, start_idx, total)
    slide_options(prs, start_idx + 1, total)
    slide_strategy(prs, start_idx + 2, total)
    slide_matrix(prs, start_idx + 3, total)
    prs.save(out_path)
    print(f"OK · on template -> {out_path}")


if __name__ == "__main__":
    out_dir = "PowerPoint-Exports"
    os.makedirs(out_dir, exist_ok=True)

    build_standalone(os.path.join(out_dir,
                                  "Valkeen_Integration-Optionen_MODERN.pptx"))
    template = "/Users/manumanera/Downloads/Entscheidungsvorlage Knauf.pptx"
    if os.path.exists(template):
        build_on_template(
            template,
            os.path.join(
                out_dir,
                "Entscheidungsvorlage_Knauf_mit_Integrationsfolien_MODERN.pptx",
            ),
        )
    else:
        print(f"WARN: Template nicht gefunden: {template}")
