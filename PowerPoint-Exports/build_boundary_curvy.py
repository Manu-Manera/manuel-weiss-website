"""
Valkeen · Technical Boundary (Originalgetreu)
=============================================

Ziel: Die obere Screenshot-Folie so genau wie möglich nachbauen:
  - Zwei KURVIGE Prozesslinien (oben Orange / unten Slate-Blau),
    die mit weicher 90°-Rundung vom 'MS Project Data'-Block abgehen.
  - Pills entlang der Linien sind WEISS mit dunkler Outline (kein Fill).
  - Jede Pill hat LINKS und RECHTS einen dunklen 'Stecker'-Punkt,
    der auf der Linie sitzt (wie Bahnstationen).
  - Pfeilspitze am Ende der Linie.
  - Golden-Rule-Callout oben links, mit dünnem grauen Connector in die
    obere Ecke des Grafikbereichs.

Die Kurven werden als echte Custom-Geometry-Shapes (cubicBezTo / quadBezTo)
direkt per XML erzeugt, damit sie editierbar bleiben.

Ausgabe:
  PowerPoint-Exports/Valkeen_TechnicalBoundary_Original.pptx
  PowerPoint-Exports/Valkeen_IntegrationOptionen_ORIGINAL_STYLE.pptx
     (3 Folien im Original-Look)
  PowerPoint-Exports/Entscheidungsvorlage_Knauf_ORIGINAL_STYLE.pptx
     (3 neue Folien an Valkeen-Knauf-Vorlage angehängt)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap
from lxml import etree

# ---------------------------------------------------------------------------
# Palette – Original-Tonalität, aber Valkeen-abgestimmt
# ---------------------------------------------------------------------------
NAVY        = RGBColor(0x00, 0x28, 0x56)  # Valkeen Navy
DARK_BLOCK  = RGBColor(0x1F, 0x2A, 0x3A)  # Quelle-Block (dunkler Anthrazit)
ORANGE      = RGBColor(0xF2, 0x9B, 0x2C)  # Allocations (warmer Orange)
ORANGE_DK   = RGBColor(0xE0, 0x85, 0x1B)  # Pfeilspitzen-Farbe (etwas dunkler)
BLUE        = RGBColor(0x1F, 0x4E, 0x9C)  # Full Schedule (satter Blau)
BLUE_DK     = RGBColor(0x19, 0x3F, 0x80)
KEY_YELLOW  = RGBColor(0xE6, 0xA8, 0x1C)
BORDER      = RGBColor(0xD0, 0xD4, 0xDA)
BORDER_DK   = RGBColor(0x4F, 0x58, 0x64)
TEXT_DARK   = RGBColor(0x1F, 0x2A, 0x3A)
TEXT_MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG          = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Calibri"
BODY_FONT  = "Calibri"

SLIDE_W, SLIDE_H = 13.333, 7.5

EMU_PER_IN = 914400
def emu(x_in: float) -> int: return int(round(x_in * EMU_PER_IN))


# ---------------------------------------------------------------------------
# Basic helpers (fill, line, text, shadow)
# ---------------------------------------------------------------------------
def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def _no_fill(shape):
    shape.fill.background()

def _line(shape, color, width_pt=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)

def _no_line(shape):
    shape.line.fill.background()

def _kill_shadow(shape):
    spPr = shape._element.spPr
    for tag in ('a:effectLst', 'a:effectRef'):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))

def _text(shape, text, *, size=12, bold=False, color=TEXT_DARK,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
          font=BODY_FONT, italic=False, line_spacing=None):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = anchor
    tf.clear()
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color

def add_rect(slide, x, y, w, h, *, fill=WHITE, line=BORDER, line_w=0.75,
             rounded=False, radius=0.04):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    if rounded:
        s.adjustments[0] = radius
    if fill is None: _no_fill(s)
    else:            _fill(s, fill)
    if line is None: _no_line(s)
    else:            _line(s, line, line_w)
    _kill_shadow(s)
    return s

def add_text(slide, x, y, w, h, text, *, size=12, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=BODY_FONT,
             italic=False, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _text(tb, text, size=size, bold=bold, color=color, align=align,
          anchor=anchor, font=font, italic=italic, line_spacing=line_spacing)
    return tb

def add_oval(slide, cx, cy, *, diam=0.18, fill=DARK_BLOCK, line=None):
    """Dunkler 'Stecker'-Punkt – wird zentriert um (cx,cy) platziert."""
    r = diam / 2
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(cx - r), Inches(cy - r),
                               Inches(diam), Inches(diam))
    _fill(s, fill)
    if line is None: _no_line(s)
    else:            _line(s, line, 0.75)
    _kill_shadow(s)
    return s

def set_slide_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


# ---------------------------------------------------------------------------
# CURVY-ARROW  (Custom Geometry via XML)
# ---------------------------------------------------------------------------
def _hex(c: RGBColor) -> str:
    return "%02X%02X%02X" % (c[0], c[1], c[2])

def _next_shape_id(slide) -> int:
    ids = []
    for el in slide.shapes._spTree.iter():
        v = el.get('id')
        if v and v.isdigit():
            ids.append(int(v))
    return (max(ids) if ids else 100) + 1

def add_curvy_stroke(slide, ops, *, color, width_pt=15.0,
                     arrow_end=True, name="CurvyStroke"):
    """
    Erzeugt eine Freeform-Stroke-Linie (nicht gefüllt) mit Bezier-Kurven.

    ops: Liste von Tuples:
        ('M', x_in, y_in)
        ('L', x_in, y_in)
        ('Q', cx_in, cy_in, x_in, y_in)
        ('C', cx1, cy1, cx2, cy2, x_in, y_in)

    Alle Koordinaten sind absolute Inches auf der Folie.
    """
    shape_id = _next_shape_id(slide)

    W = emu(SLIDE_W)
    H = emu(SLIDE_H)

    path_ops = []
    for op in ops:
        tag = op[0]
        if tag == 'M':
            path_ops.append(
                f'<a:moveTo><a:pt x="{emu(op[1])}" y="{emu(op[2])}"/></a:moveTo>')
        elif tag == 'L':
            path_ops.append(
                f'<a:lnTo><a:pt x="{emu(op[1])}" y="{emu(op[2])}"/></a:lnTo>')
        elif tag == 'Q':
            path_ops.append(
                '<a:quadBezTo>'
                f'<a:pt x="{emu(op[1])}" y="{emu(op[2])}"/>'
                f'<a:pt x="{emu(op[3])}" y="{emu(op[4])}"/>'
                '</a:quadBezTo>')
        elif tag == 'C':
            path_ops.append(
                '<a:cubicBezTo>'
                f'<a:pt x="{emu(op[1])}" y="{emu(op[2])}"/>'
                f'<a:pt x="{emu(op[3])}" y="{emu(op[4])}"/>'
                f'<a:pt x="{emu(op[5])}" y="{emu(op[6])}"/>'
                '</a:cubicBezTo>')

    tailEnd = (
        '<a:tailEnd type="triangle" w="lg" h="lg"/>'
        if arrow_end else '')

    stroke_emu = int(width_pt * 12700)
    hex_color = _hex(color)

    xml = f'''<p:sp xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
                    xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                    xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <p:nvSpPr>
    <p:cNvPr id="{shape_id}" name="{name}"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="0" y="0"/><a:ext cx="{W}" cy="{H}"/></a:xfrm>
    <a:custGeom>
      <a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>
      <a:rect l="0" t="0" r="{W}" b="{H}"/>
      <a:pathLst>
        <a:path w="{W}" h="{H}">
          {"".join(path_ops)}
        </a:path>
      </a:pathLst>
    </a:custGeom>
    <a:noFill/>
    <a:ln w="{stroke_emu}" cap="rnd" cmpd="sng" algn="ctr">
      <a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>
      <a:prstDash val="solid"/>
      <a:round/>
      {tailEnd}
    </a:ln>
    <a:effectLst/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr/>
    <a:lstStyle/>
    <a:p><a:endParaRPr lang="de-DE"/></a:p>
  </p:txBody>
</p:sp>'''

    sp_el = etree.fromstring(xml)
    slide.shapes._spTree.append(sp_el)
    return sp_el


# ---------------------------------------------------------------------------
# Pill (Rounded Rect mit Stecker-Punkten an den Enden)
# ---------------------------------------------------------------------------
def add_station_pill(slide, cx, cy, *, w=1.50, h=0.60, text,
                     border=DARK_BLOCK, border_w_pt=2.0,
                     fill=WHITE, text_color=TEXT_DARK, size=10.5, bold=True,
                     text_lines=1, dot_diam=0.17):
    """
    Eine 'Bahnstations-Pill' – zentriert an (cx, cy) auf der Prozesslinie.
    Enthält zwei dunkle Stecker-Punkte links und rechts (auf der Linie).
    """
    x = cx - w / 2
    y = cy - h / 2
    # Pill
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.adjustments[0] = 0.5
    _fill(s, fill)
    _line(s, border, border_w_pt)
    _kill_shadow(s)
    # Text
    _text(s, text, size=size, bold=bold, color=text_color,
          font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
          line_spacing=1.05)
    # Stecker-Punkte – zentriert über dem Pill-Rand
    add_oval(slide, cx - w / 2, cy, diam=dot_diam, fill=border)
    add_oval(slide, cx + w / 2, cy, diam=dot_diam, fill=border)
    return s


# ---------------------------------------------------------------------------
# TITLE (wie Screenshot: Bold Dunkel + Bold Grau)
# ---------------------------------------------------------------------------
def add_title_two_tone(slide, bold_black, bold_grey, *, y=0.30):
    tb = slide.shapes.add_textbox(Inches(0.50), Inches(y),
                                  Inches(12.3), Inches(0.85))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = bold_black
    r1.font.name = TITLE_FONT
    r1.font.size = Pt(30)
    r1.font.bold = True
    r1.font.color.rgb = TEXT_DARK
    r2 = p.add_run()
    r2.text = bold_grey
    r2.font.name = TITLE_FONT
    r2.font.size = Pt(30)
    r2.font.bold = True
    r2.font.color.rgb = RGBColor(0x88, 0x90, 0x9A)


# ===========================================================================
# SLIDE 1  –  Technical Boundary (Originalgetreu, kurvige Prozesslinien)
# ===========================================================================
def slide_boundary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)

    # --- Titel ----------------------------------------------------------
    add_title_two_tone(slide,
                       "The Technical Boundary: ",
                       "Schedule vs. Allocation")

    # === Geometrie (sorgfältig auf Original abgestimmt) =================
    # MS Project Data – Source (links)
    src_x, src_y, src_w, src_h = 0.60, 3.55, 2.10, 1.20
    src_mid_y = src_y + src_h / 2    # 4.15

    # Orange Prozesslinie (oben)
    y_top_line   = 2.95        # horizontal-Niveau der orange Linie
    y_top_src    = src_y + 0.35   # Austrittspunkt oben am Source-Block
    # Blaue Prozesslinie (unten)
    y_bot_line   = 5.30
    y_bot_src    = src_y + src_h - 0.35  # Austrittspunkt unten am Source-Block

    x_bend       = src_x + src_w + 0.55   # x-Position der runden Ecke
    x_line_start = x_bend + 0.25          # Beginn des horizontalen Teils

    # Zielblöcke rechts
    tgt_x = 10.95
    tgt_w = 1.95
    tgt_h = 0.85
    x_arrow_end = tgt_x - 0.02            # Pfeilspitze direkt an den Block

    # --- Golden Rule Callout (oben links) ------------------------------
    cal_x, cal_y, cal_w, cal_h = 0.60, 1.30, 5.90, 0.95
    add_rect(slide, cal_x, cal_y, cal_w, cal_h, fill=WHITE, line=BORDER,
             line_w=0.75, rounded=True, radius=0.03)
    # Key-Icon-Hintergrund (Unicode-Schlüssel in Gold)
    key_size = 0.58
    key_cx = cal_x + 0.45
    key_cy = cal_y + cal_h / 2
    k = slide.shapes.add_textbox(Inches(key_cx - key_size / 2),
                                 Inches(key_cy - key_size / 2),
                                 Inches(key_size), Inches(key_size))
    _text(k, "🔑", size=26, color=KEY_YELLOW, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE, bold=False)
    # Text rechts
    add_text(slide, cal_x + 1.00, cal_y + 0.11, cal_w - 1.10, cal_h - 0.22,
             "The Golden Rule:  The destination of your data\n"
             "dictates your integration tool.",
             size=11.5, color=TEXT_DARK, font=BODY_FONT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.25)
    # Callout-Connector (dünne graue L-Linie nach rechts oben)
    add_curvy_stroke(
        slide,
        ops=[
            ('M', cal_x + cal_w, cal_y + cal_h / 2),
            ('L', cal_x + cal_w + 0.25, cal_y + cal_h / 2),
            ('Q',
             cal_x + cal_w + 0.45, cal_y + cal_h / 2,
             cal_x + cal_w + 0.45, cal_y + cal_h / 2 - 0.20),
            ('L', cal_x + cal_w + 0.45, 1.10),
            ('L', 12.75, 1.10),
            ('L', 12.75, y_top_line - 0.30),
        ],
        color=BORDER, width_pt=0.75, arrow_end=False,
        name="CalloutConnector")

    # --- MS Project Data Block -----------------------------------------
    src = add_rect(slide, src_x, src_y, src_w, src_h,
                   fill=DARK_BLOCK, line=None, rounded=True, radius=0.06)
    _text(src, "MS Project Data", size=18, bold=True, color=WHITE,
          font=TITLE_FONT, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    # =========================================================================
    # ORANGE Prozesslinie (oben) – kurvig
    #    Start: rechte Kante Source, obere Hälfte
    #    Ende:  linker Rand Allocations-Block
    # =========================================================================
    x_src_right = src_x + src_w  # 2.70
    add_curvy_stroke(
        slide,
        ops=[
            ('M', x_src_right, y_top_src),
            ('L', x_bend - 0.25, y_top_src),
            # runde Ecke (nach oben biegen)
            ('Q',
             x_bend,        y_top_src,
             x_bend,        y_top_src - 0.25),
            ('L', x_bend,    y_top_line + 0.25),
            # runde Ecke (nach rechts biegen)
            ('Q',
             x_bend,        y_top_line,
             x_bend + 0.25, y_top_line),
            # lange horizontale Strecke
            ('L', x_arrow_end, y_top_line),
        ],
        color=ORANGE, width_pt=15.0, arrow_end=True, name="OrangePath")

    # Start-Stecker am Source-Block (wo die Linie abgeht)
    add_oval(slide, x_src_right, y_top_src, diam=0.22,
             fill=ORANGE_DK)

    # =========================================================================
    # BLAU  Prozesslinie (unten) – kurvig
    # =========================================================================
    add_curvy_stroke(
        slide,
        ops=[
            ('M', x_src_right, y_bot_src),
            ('L', x_bend - 0.25, y_bot_src),
            ('Q',
             x_bend,        y_bot_src,
             x_bend,        y_bot_src + 0.25),
            ('L', x_bend,    y_bot_line - 0.25),
            ('Q',
             x_bend,        y_bot_line,
             x_bend + 0.25, y_bot_line),
            ('L', x_arrow_end, y_bot_line),
        ],
        color=BLUE, width_pt=15.0, arrow_end=True, name="BluePath")
    add_oval(slide, x_src_right, y_bot_src, diam=0.22, fill=BLUE_DK)

    # =========================================================================
    # PILLS auf der ORANGE-Linie (3 Stück, äquidistant)
    # =========================================================================
    # Nutzbarer Bereich horizontal: von x_line_start bis x_arrow_end - 0.2
    pill_area_x0 = x_line_start + 0.10      # etwas Abstand zur Kurve
    pill_area_x1 = x_arrow_end - 0.10

    pill_h_std = 0.62
    pill_h_tall = 0.80   # für 2-zeilige Pill "Project Online Sync"

    orange_pills = [
        {"text": "Project Online\nSync", "w": 1.75, "h": pill_h_tall,
         "lines": 2},
        {"text": "MSP Add-In",          "w": 1.50, "h": pill_h_std,
         "lines": 1},
        {"text": "MPP Import",          "w": 1.50, "h": pill_h_std,
         "lines": 1},
    ]

    # Gleichmässig verteilen
    n = len(orange_pills)
    span = pill_area_x1 - pill_area_x0
    for i, p in enumerate(orange_pills):
        cx = pill_area_x0 + span * (i + 0.5) / n
        add_station_pill(slide, cx, y_top_line,
                         w=p["w"], h=p["h"], text=p["text"],
                         border=DARK_BLOCK, border_w_pt=2.0,
                         fill=WHITE, text_color=TEXT_DARK,
                         size=10.5, bold=True, dot_diam=0.17)

    # =========================================================================
    # PILL auf der BLAU-Linie (1 Stück, mittig)
    # =========================================================================
    blue_cx = (pill_area_x0 + pill_area_x1) / 2
    add_station_pill(slide, blue_cx, y_bot_line,
                     w=1.50, h=pill_h_std, text="MPP Import",
                     border=DARK_BLOCK, border_w_pt=2.0,
                     fill=WHITE, text_color=TEXT_DARK,
                     size=10.5, bold=True, dot_diam=0.17)

    # =========================================================================
    # Zielblöcke rechts (Allocations Grid / Full Schedule)
    # =========================================================================
    # Allocations Grid – orange gefüllt
    tgt_top_y = y_top_line - tgt_h / 2
    t1 = add_rect(slide, tgt_x, tgt_top_y, tgt_w, tgt_h,
                  fill=ORANGE, line=None, rounded=True, radius=0.04)
    _text(t1, "Allocations\nGrid", size=14, bold=True, color=WHITE,
          font=TITLE_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
          line_spacing=1.05)
    # Beschreibung darunter
    add_text(slide, tgt_x, tgt_top_y + tgt_h + 0.10, tgt_w + 0.30, 0.95,
             "High-level planned and actual\n"
             "resource demand. Supported by\n"
             "all synchronisation methods.",
             size=8.5, color=TEXT_MUTED, font=BODY_FONT,
             align=PP_ALIGN.LEFT, line_spacing=1.25)

    # Full Schedule – blau gefüllt
    tgt_bot_y = y_bot_line - tgt_h / 2
    t2 = add_rect(slide, tgt_x, tgt_bot_y, tgt_w, tgt_h,
                  fill=BLUE, line=None, rounded=True, radius=0.04)
    _text(t2, "Full Schedule", size=14, bold=True, color=WHITE,
          font=TITLE_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, tgt_x, tgt_bot_y + tgt_h + 0.10, tgt_w + 0.30, 1.00,
             "Granular task logic, dependencies,\n"
             "task types, and dates.\n"
             "Supported only by MPP Import.",
             size=8.5, color=TEXT_MUTED, font=BODY_FONT,
             align=PP_ALIGN.LEFT, line_spacing=1.25)


# ===========================================================================
# SLIDE 2  –  Integration Options: At a Glance (Originalgetreu)
# ===========================================================================
def slide_options(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)

    add_title_two_tone(slide, "Integration Options: ", "At a Glance")

    # Table
    left = 0.60
    top  = 1.25
    lbl_w = 1.40
    col_w = (13.33 - 2 * left - lbl_w) / 3
    total_w = lbl_w + 3 * col_w

    # ---- Header (Tool names) -----------------------------------------
    headers = ["Project Online Sync", "MSP Add-In", "MPP Import"]
    head_h  = 0.75
    for i, h in enumerate(headers):
        x = left + lbl_w + i * col_w
        add_text(slide, x, top + 0.10, col_w, 0.55, h,
                 size=18, bold=True, color=TEXT_DARK, font=TITLE_FONT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Horizontale Trennlinien
    def hsep(y, color=RGBColor(0xBD, 0xC2, 0xCA), w_pt=0.75):
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        Inches(left), Inches(y),
                                        Inches(left + total_w), Inches(y))
        ln.line.color.rgb = color
        ln.line.width = Pt(w_pt)

    # Vertikale Trennlinien
    def vsep(x, y1, y2, color=RGBColor(0xBD, 0xC2, 0xCA), w_pt=0.75):
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        Inches(x), Inches(y1),
                                        Inches(x), Inches(y2))
        ln.line.color.rgb = color
        ln.line.width = Pt(w_pt)

    hsep(top + head_h)
    # vertical separators
    for i in range(4):
        vsep(left + lbl_w + i * col_w, top + head_h, top + 6.0)

    # ---- Row 1: Data Target ------------------------------------------
    row1_y = top + head_h
    row1_h = 0.95
    # Row-Label (links)
    add_text(slide, left + 0.10, row1_y + 0.15, lbl_w - 0.20, row1_h - 0.30,
             "Data\nTarget", size=13, bold=True, color=TEXT_DARK,
             font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
             line_spacing=1.1)

    def center_pill(cx, cy, w, h, text, fill, text_color=WHITE, size=11):
        x = cx - w / 2
        y = cy - h / 2
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y),
                                   Inches(w), Inches(h))
        s.adjustments[0] = 0.5
        _fill(s, fill)
        _no_line(s)
        _kill_shadow(s)
        _text(s, text, size=size, bold=True, color=text_color,
              font=BODY_FONT)
        return s

    # Spalte 1 & 2 – "Allocations ONLY" (orange)
    for i in (0, 1):
        cx = left + lbl_w + i * col_w + col_w / 2
        cy = row1_y + row1_h / 2
        center_pill(cx, cy, 2.30, 0.48, "Allocations ONLY", ORANGE,
                    text_color=WHITE, size=11)
    # Spalte 3 – Full Schedule (blau) + Allocations (orange) nebeneinander
    cx = left + lbl_w + 2 * col_w + col_w / 2
    cy = row1_y + row1_h / 2
    # zwei Pills nebeneinander
    g = 0.15
    w1, w2 = 1.70, 1.55
    center_pill(cx - (w2 + g) / 2, cy, w1, 0.48, "Full Schedule", BLUE,
                text_color=WHITE, size=11)
    center_pill(cx + (w1 + g) / 2, cy, w2, 0.48, "Allocations", ORANGE,
                text_color=WHITE, size=11)

    hsep(row1_y + row1_h)

    # ---- Row 2: Pros -------------------------------------------------
    row2_y = row1_y + row1_h
    row2_h = 2.55
    add_text(slide, left + 0.10, row2_y + 0.20, lbl_w - 0.20, 0.40, "Pros",
             size=13, bold=True, color=TEXT_DARK, font=BODY_FONT,
             align=PP_ALIGN.LEFT)

    pros = [
        ["Enables automated, scheduled synchronisation",
         "Admin-led governance ensures systemic control",
         "Centralised process"],
        ["Complete Project Manager control",
         "Features a pre-sync data validation\nand error preview interface",
         "Enables advanced WBS mapping"],
        ["The only tool that imports full\nschedules (tasks, dependencies,\ndates, and durations)",
         "Excellent for initial system migrations"],
    ]
    for i, items in enumerate(pros):
        x = left + lbl_w + i * col_w
        tb = slide.shapes.add_textbox(Inches(x + 0.20), Inches(row2_y + 0.20),
                                      Inches(col_w - 0.30),
                                      Inches(row2_h - 0.40))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        for k, line in enumerate(items):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(8)
            # Bulletpoint als Strich mit Semi-Bold-Gewicht
            r1 = p.add_run()
            r1.text = "- "
            r1.font.size = Pt(11.5)
            r1.font.color.rgb = TEXT_DARK
            r1.font.name = BODY_FONT
            r2 = p.add_run()
            r2.text = line
            r2.font.size = Pt(11)
            r2.font.color.rgb = TEXT_DARK
            r2.font.name = BODY_FONT

    hsep(row2_y + row2_h)

    # ---- Row 3: Cons -------------------------------------------------
    row3_y = row2_y + row2_h
    row3_h = 2.0
    add_text(slide, left + 0.10, row3_y + 0.20, lbl_w - 0.20, 0.40, "Cons",
             size=13, bold=True, color=TEXT_DARK, font=BODY_FONT,
             align=PP_ALIGN.LEFT)

    cons = [
        ["Physically incapable of importing\nschedule data",
         "Requires root URL and global API\ncredential setup"],
        ["Requires manual execution per project",
         "Requires individual user API token\ngeneration"],
        ["A manual, file-based process\nrequiring drag-and-drop execution",
         "Schedule imports automatically\noverwrite existing schedule entities"],
    ]
    for i, items in enumerate(cons):
        x = left + lbl_w + i * col_w
        tb = slide.shapes.add_textbox(Inches(x + 0.20), Inches(row3_y + 0.20),
                                      Inches(col_w - 0.30),
                                      Inches(row3_h - 0.40))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        for k, line in enumerate(items):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(8)
            # Warn-Icon mittels Kreisform
            r1 = p.add_run()
            r1.text = "! "
            r1.font.size = Pt(11.5)
            r1.font.bold = True
            r1.font.color.rgb = ORANGE_DK
            r1.font.name = BODY_FONT
            r2 = p.add_run()
            r2.text = line
            r2.font.size = Pt(11)
            r2.font.color.rgb = TEXT_DARK
            r2.font.name = BODY_FONT

    hsep(row3_y + row3_h)


# ===========================================================================
# SLIDE 3  –  Selecting Your Integration Strategy (3 Karten, originalgetreu)
# ===========================================================================
def slide_strategy(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)

    add_title_two_tone(slide, "Selecting Your ", "Integration Strategy")

    cards = [
        {
            "icon": "converge",
            "context": "If your operational model is…",
            "body": "Centralised planning requiring automated, hands-off "
                    "synchronisation across the enterprise.",
            "deploy": "Then deploy",
            "tool": "Project\nOnline Sync",
        },
        {
            "icon": "network",
            "context": "If your operational model is…",
            "body": "Decentralised planning where individual Project Managers "
                    "must control, preview, and validate their own data "
                    "before submission.",
            "deploy": "Then deploy the",
            "tool": "MSP\nAdd-In",
        },
        {
            "icon": "disk",
            "context": "If your operational model requires…",
            "body": "A one-time system migration, or an ongoing necessity to "
                    "view full task logic, dependencies, and granular "
                    "project schedules in Tempus.",
            "deploy": "Then utilise",
            "tool": "MPP\nImport",
        },
    ]

    card_y = 1.25
    card_w = 4.05
    card_h = 5.85
    gap = 0.15
    start_x = 0.50 + (13.333 - 1.0 - 3 * card_w - 2 * gap) / 2

    for i, c in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rect(slide, x, card_y, card_w, card_h,
                 fill=WHITE, line=RGBColor(0xC7, 0xCD, 0xD3), line_w=0.75,
                 rounded=True, radius=0.015)

        icon_cx = x + card_w / 2
        icon_cy = card_y + 1.15

        if c["icon"] == "converge":
            _icon_converge(slide, icon_cx, icon_cy)
        elif c["icon"] == "network":
            _icon_network(slide, icon_cx, icon_cy)
        elif c["icon"] == "disk":
            _icon_disk(slide, icon_cx, icon_cy)

        # Context
        add_text(slide, x + 0.30, card_y + 2.45, card_w - 0.60, 0.35,
                 c["context"], size=11.5, bold=True, color=TEXT_DARK,
                 font=BODY_FONT, align=PP_ALIGN.CENTER)
        # Body
        add_text(slide, x + 0.35, card_y + 2.85, card_w - 0.70, 1.40,
                 c["body"], size=11, color=TEXT_MUTED, font=BODY_FONT,
                 align=PP_ALIGN.CENTER, line_spacing=1.3)

        # Down-Arrow
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     Inches(icon_cx - 0.22),
                                     Inches(card_y + 4.25),
                                     Inches(0.44), Inches(0.45))
        _fill(arr, DARK_BLOCK); _no_line(arr); _kill_shadow(arr)

        # Then deploy
        add_text(slide, x + 0.30, card_y + 4.78, card_w - 0.60, 0.30,
                 c["deploy"], size=11, bold=True, color=TEXT_DARK,
                 font=BODY_FONT, align=PP_ALIGN.CENTER)
        # Tool-Name sehr gross, bold
        add_text(slide, x + 0.20, card_y + 5.05, card_w - 0.40, 0.85,
                 c["tool"], size=22, bold=True, color=TEXT_DARK,
                 font=TITLE_FONT, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.0)


def _icon_converge(slide, cx, cy):
    """Zentraler Kreis + 4 Pfeile die nach innen zeigen (Karte 1)."""
    r = 0.28
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(cx - r), Inches(cy - r),
                               Inches(r * 2), Inches(r * 2))
    _fill(c, DARK_BLOCK); _no_line(c); _kill_shadow(c)

    # 4 Pfeile: Bounding-Box ist alen (lang) × aw (breit).
    # Durch s.rotation um 90/180/270 wird die Bounding-Box IN-PLACE rotiert
    # (Shape dreht um eigenes Center). Deshalb platzieren wir die Box so,
    # dass deren Center an der gewünschten Position sitzt.
    alen, aw = 0.55, 0.26
    gap = 0.18   # Abstand von der Kreiskante zur Pfeilspitze

    def _arrow(rot, tip_cx, tip_cy):
        """
        Platziert einen RIGHT_ARROW, dessen SPITZE-Center an (tip_cx, tip_cy)
        zeigt (nach Rotation).
        """
        # Unrotated Bounding-Box-Center = (box.x + alen/2, box.y + aw/2).
        # Wir wollen: Center nach Rotation = (tip_cx - alen/2 * dir, tip_cy).
        # Einfacher: Center der Shape = gewünschtes Shape-Center, dann
        # Spitze ergibt sich durch Rotation automatisch.
        # Shape-Center soll so liegen, dass die Spitze (Rotation beachtet)
        # exakt gap von (cx,cy) entfernt ist.
        import math
        # Richtung vom Zentrum zur Spitze
        rad = math.radians(rot)
        # RIGHT_ARROW (rot=0) hat Spitze rechts, also Richtung = (cos, sin)
        dx, dy = math.cos(rad), math.sin(rad)
        # Shape-Center = Zentrum_der_Ziel_Spitze - (alen/2) * Richtung
        tip_x = cx + (r + gap + alen) * dx  # ursprünglich gewollte Spitze
        tip_y = cy + (r + gap + alen) * dy
        # ... aber wir wollen die Spitze GENAU bei (r+gap) vom Center,
        # und die Shape so drehen, dass sie nach innen zeigt.
        # Einfacher praktischer Ansatz:
        # Place bounding box with its center at distance (r + gap + alen/2)
        # from (cx,cy), and rotate so pointing inward.
        center_x = cx + (r + gap + alen / 2) * dx
        center_y = cy + (r + gap + alen / 2) * dy
        # Rotation soll die Spitze nach innen zeigen = entgegengesetzt zur
        # Richtung vom Zentrum aus betrachtet.
        # RIGHT_ARROW zeigt nach rechts bei rot=0.
        # Wenn Pfeil rechts vom Zentrum sitzt (dx=+1), soll Spitze nach links
        # zeigen → rotation=180°.
        rot_angle = (rot + 180) % 360
        s = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(center_x - alen / 2), Inches(center_y - aw / 2),
            Inches(alen), Inches(aw))
        s.rotation = rot_angle
        _fill(s, DARK_BLOCK); _no_line(s); _kill_shadow(s)
        try: s.adjustments[0] = 0.55; s.adjustments[1] = 0.35
        except Exception: pass

    # 4 Pfeile – Osten / Westen / Norden / Süden vom Kreis aus gesehen
    _arrow(0,   None, None)    # Pfeil rechts vom Kreis → zeigt nach links
    _arrow(180, None, None)    # Pfeil links vom Kreis → zeigt nach rechts
    _arrow(90,  None, None)    # Pfeil unten vom Kreis → zeigt nach oben
    _arrow(270, None, None)    # Pfeil oben vom Kreis → zeigt nach unten


def _icon_network(slide, cx, cy):
    """Hub-and-Spoke: zentraler Knoten + 6 kleine Personen-Knoten (Karte 2)."""
    import math
    big_r = 0.28
    b = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(cx - big_r), Inches(cy - big_r),
                               Inches(big_r * 2), Inches(big_r * 2))
    _fill(b, DARK_BLOCK); _no_line(b); _kill_shadow(b)

    R = 0.95
    hr = 0.16  # head radius
    for k in range(6):
        ang = math.radians(60 * k - 90)
        sx = cx + R * math.cos(ang)
        sy = cy + R * math.sin(ang)
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        Inches(cx), Inches(cy),
                                        Inches(sx), Inches(sy))
        ln.line.color.rgb = DARK_BLOCK
        ln.line.width = Pt(1.25)
        # Kopf
        h = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(sx - hr), Inches(sy - hr - 0.03),
                                   Inches(hr * 2), Inches(hr * 2))
        _fill(h, DARK_BLOCK); _no_line(h); _kill_shadow(h)
        # Körper
        bw, bh = hr * 2.2, hr * 1.05
        bd = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(sx - bw / 2),
                                    Inches(sy + hr - 0.02),
                                    Inches(bw), Inches(bh))
        _fill(bd, DARK_BLOCK); _no_line(bd); _kill_shadow(bd)


def _icon_disk(slide, cx, cy):
    """HDD-Stack mit Migrations-Pfeil nach rechts (Karte 3).

    Umsetzung als saubere Cylinder-Shape (MSO_SHAPE.CAN) links vom Icon
    und ein gebogener Pfeil rechts daneben. Icon leicht nach links verschoben,
    damit der Pfeil Platz hat.
    """
    # Cylinder (HDD)
    cyl_w = 1.10
    cyl_h = 1.25
    cyl_x = cx - 0.95
    cyl_y = cy - cyl_h / 2
    cyl = slide.shapes.add_shape(MSO_SHAPE.CAN,
                                 Inches(cyl_x), Inches(cyl_y),
                                 Inches(cyl_w), Inches(cyl_h))
    cyl.rotation = 90  # CAN ist vertikal; 90° -> horizontal wäre Scheibenstapel von der Seite
    # Zurücksetzen – vertikal ist OK wie ein HDD-Stack
    cyl.rotation = 0
    _fill(cyl, DARK_BLOCK); _no_line(cyl); _kill_shadow(cyl)

    # Gebogener Pfeil rechts daneben
    arr = slide.shapes.add_shape(MSO_SHAPE.CURVED_RIGHT_ARROW,
                                 Inches(cx + 0.05), Inches(cy - 0.85),
                                 Inches(0.85), Inches(1.60))
    _fill(arr, DARK_BLOCK); _no_line(arr); _kill_shadow(arr)


# ===========================================================================
# Build
# ===========================================================================
def build_standalone(out_path, *, only_boundary=False):
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide_boundary(prs)
    if not only_boundary:
        slide_options(prs)
        slide_strategy(prs)
    prs.save(out_path)
    print(f"OK  ->  {out_path}")


def build_on_template(template_path, out_path):
    prs = Presentation(template_path)
    slide_boundary(prs)
    slide_options(prs)
    slide_strategy(prs)
    prs.save(out_path)
    print(f"OK  ->  {out_path}")


if __name__ == "__main__":
    out_dir = "PowerPoint-Exports"
    os.makedirs(out_dir, exist_ok=True)

    # Nur die Boundary-Folie für schnelles Review
    build_standalone(
        os.path.join(out_dir, "Valkeen_TechnicalBoundary_ORIGINAL.pptx"),
        only_boundary=True,
    )

    # Komplettes Deck mit allen drei Folien im Original-Look
    build_standalone(
        os.path.join(out_dir, "Valkeen_IntegrationOptionen_ORIGINAL_STYLE.pptx"),
    )

    template = "/Users/manumanera/Downloads/Entscheidungsvorlage Knauf.pptx"
    if os.path.exists(template):
        build_on_template(
            template,
            os.path.join(
                out_dir,
                "Entscheidungsvorlage_Knauf_ORIGINAL_STYLE.pptx",
            ),
        )
    else:
        print(f"WARN: Template nicht gefunden: {template}")
