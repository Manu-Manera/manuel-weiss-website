#!/usr/bin/env python3
"""
Baut "Tempus Excel-Upload Walkthrough" als PPTX EXAKT im Valkeen/Tempus-Master-Design.

Methode: Klont echte Folien aus valkeen-2026-master.pptx (Theme "Blau II", Calibri,
Tempus-Branding/Logo) und ersetzt nur den Text. Kein Neu-Design — identischer Look.

Vorlagen-Folien im Master:
  - Slide 0  : Titelfolie (Panel links, Titel, Logos)
  - Slide 50 : Agenda (Panel + Titel + Liste)
  - Slide 26 : Inhaltsfolie "Default Slide" (Titel + Eyebrow + Body + Logo-Gruppe)
"""

from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

MASTER = Path("/Users/manumanera/Downloads/valkeen-2026-master.pptx")
OUT = Path(
    "/Users/manumanera/Documents/GitHub/Persönliche Website/"
    "Onboarding Valkeen/onboarding-app/public/tempus-excel-upload-walkthrough.pptx"
)

TPL_TITLE = 0
TPL_AGENDA = 50
TPL_CONTENT = 26


# ── Slide-Cloning ───────────────────────────────────────────────
_RID_ATTRS = (
    qn("r:embed"),
    qn("r:link"),
    qn("r:id"),
    qn("r:pict"),
    qn("r:dm"),
    qn("r:lo"),
    qn("r:qs"),
    qn("r:cs"),
)


def clone_slide(prs, src_index):
    """Dupliziert eine vorhandene Folie inkl. Bild-Relationships ans Ende."""
    src = prs.slides[src_index]
    layout = src.slide_layout
    dst = prs.slides.add_slide(layout)
    # vom Layout automatisch erzeugte Platzhalter entfernen
    for sh in list(dst.shapes):
        sh._element.getparent().remove(sh._element)

    # Relationships übernehmen → Mapping alte rId → neue rId
    rid_map = {}
    for rid, rel in src.part.rels.items():
        if rel.reltype.endswith("/slideLayout"):
            continue
        if rel.is_external:
            new = dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new = dst.part.rels.get_or_add(rel.reltype, rel._target)
        rid_map[rid] = new

    # Shapes kopieren und rId-Referenzen remappen
    for sh in src.shapes:
        el = copy.deepcopy(sh._element)
        for node in [el] + el.findall(".//*"):
            for attr in _RID_ATTRS:
                val = node.get(attr)
                if val and val in rid_map:
                    node.set(attr, rid_map[val])
        dst.shapes._spTree.append(el)
    return dst


def delete_slides(prs, indices):
    """Entfernt Folien per Index (und droppt deren Relationships)."""
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    for i in sorted(indices, reverse=True):
        sldId = ids[i]
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def cleanup_orphan_slide_parts(pptx_path: Path) -> int:
    """
    Entfernt verwaiste ppt/slides/*-Parts aus der ZIP.
    python-pptx droppt Relationships, lässt aber alte Slide-XML-Dateien stehen → PowerPoint-Reparatur.
    """
    import re
    import zipfile
    from io import BytesIO
    from xml.etree import ElementTree as ET

    CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"

    with zipfile.ZipFile(pptx_path, "r") as zin:
        pres_rels = zin.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
        kept_slides = set(re.findall(r"slides/(slide\d+\.xml)", pres_rels))
        if not kept_slides:
            return 0

        removed = 0
        buf = BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
            ct_xml = zin.read("[Content_Types].xml")
            ct_root = ET.fromstring(ct_xml)
            ct_to_remove = []

            for item in zin.infolist():
                name = item.filename
                base = name.split("/")[-1]

                if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                    if base not in kept_slides:
                        removed += 1
                        ct_to_remove.append(f"/ppt/slides/{base}")
                        continue
                if name.startswith("ppt/slides/_rels/slide") and name.endswith(".xml.rels"):
                    slide_xml = base.replace(".xml.rels", ".xml")
                    if slide_xml not in kept_slides:
                        removed += 1
                        continue

                if name == "[Content_Types].xml":
                    for ov in list(ct_root):
                        if ov.tag.endswith("Override"):
                            part = ov.get("PartName", "")
                            if part in ct_to_remove:
                                ct_root.remove(ov)
                    ct_xml = ET.tostring(ct_root, encoding="utf-8", xml_declaration=True)
                    zout.writestr(item, ct_xml)
                else:
                    zout.writestr(item, zin.read(name))

    if removed:
        pptx_path.write_bytes(buf.getvalue())
    return removed


# ── Text-Helfer ─────────────────────────────────────────────────
def _shapes_with_text(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def set_title(shape, text):
    tf = shape.text_frame
    # erste Run-Formatierung merken
    p0 = tf.paragraphs[0]
    run = p0.runs[0] if p0.runs else p0.add_run()
    run.text = text
    # zusätzliche Absätze leeren
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    for r in p0.runs[1:]:
        r._r.getparent().remove(r._r)


def set_body(shape, items):
    """items: Liste von (text, level) – behält Font/Farbe des Originals."""
    tf = shape.text_frame
    base = tf.paragraphs[0]
    base_run = base.runs[0] if base.runs else None
    sz = base_run.font.size if base_run else None
    bold = base_run.font.bold if base_run else None
    color_rgb = None
    try:
        if base_run and base_run.font.color and base_run.font.color.type is not None:
            color_rgb = base_run.font.color.rgb
    except Exception:
        color_rgb = None
    name = base_run.font.name if base_run else None

    # alle Absätze entfernen
    for p in list(tf.paragraphs):
        p._p.getparent().remove(p._p)

    for text, level in items:
        p = tf.add_paragraph()
        p.level = max(0, min(level, 4))
        r = p.add_run()
        r.text = text
        if sz is not None:
            r.font.size = sz
        if name:
            r.font.name = name
        if bold is not None and level == 0 and text.endswith(":"):
            r.font.bold = True
        if color_rgb is not None:
            try:
                r.font.color.rgb = color_rgb
            except Exception:
                pass


# ── Inhalt ──────────────────────────────────────────────────────
TITLE = "Excel-Upload Walkthrough"
SUBTITLE = "So passen die Dateien in Tempus zusammen"

AGENDA = [
    ("Das große Ganze – warum strukturierter Upload", 0),
    ("Wo der Upload passiert", 0),
    ("Die 10 Import-Dateien", 0),
    ("Zusammenhänge & Abhängigkeiten", 0),
    ("Die richtige Reihenfolge (1 → 10)", 0),
    ("Wie Tempus Zeilen erkennt", 0),
    ("Import Behavior: Merge / Overwrite / Skip", 0),
    ("Zeit-phasierte Spalten", 0),
    ("Schritt-für-Schritt & häufige Fehler", 0),
    ("Checkliste & nächste Schritte", 0),
]

# (Titel, Eyebrow, [ (bullet, level) ... ])
CONTENT = [
    (
        "Das große Ganze",
        "Tempus ist anfangs ein leeres Blatt",
        [
            ("Wir befüllen Tempus über standardisierte Excel-Vorlagen – jede steht für einen Datentyp.", 0),
            ("Stammdaten zuerst: Felder (Attribute), Ressourcen, Skills", 0),
            ("Dann die Container: Projekte (nutzen die Attribute)", 0),
            ("Dann die Verknüpfungen: Zuweisungen = Ressource + Projekt + Aufwand", 0),
            ("Zuletzt die Details: Finanzen, Raten, Admin-Zeiten, Team", 0),
            ("Faustregel: Felder → Ressourcen → Projekte → Zuweisungen → Details", 0),
        ],
    ),
    (
        "Wo der Upload passiert",
        "Ein Ort, viele Vorlagen",
        [
            ("Navigieren: Admin Settings > Data Synchronization > Excel", 0),
            ("Datei wählen: Choose File oder .xlsx auf die Seite ziehen", 0),
            ("Synchronisieren: Button Synchronize startet den Import", 0),
            ("Wichtig: Die auszufüllende Vorlage muss das erste Tab sein", 0),
            ("Format immer .xlsx – pro Synchronisation eine Vorlage / ein Datentyp", 0),
            ("Empfehlung: erst Sandbox testen, dann Produktion", 0),
        ],
    ),
    (
        "Die 10 Import-Dateien (1/2)",
        "Fundament & Aufbau",
        [
            ("Fundament (keine Voraussetzung):", 0),
            ("1 · Attributes — eigene Felder & Auswahllisten", 1),
            ("2 · Resources — Ressourcen-Pool, Kapazität, Rollen", 1),
            ("6 · Skills — Skill-Katalog & Kategorien", 1),
            ("Baut auf dem Fundament auf:", 0),
            ("3 · Projects — Projektliste (← Attribute)", 1),
            ("8 · Advanced Rates — Sätze je Ressource (← Ressourcen)", 1),
            ("10 · Team Resources — Team-Mitglieder (← Ressourcen)", 1),
            ("5 · Admin Times — Nicht-Projektzeiten (← Ressourcen)", 1),
        ],
    ),
    (
        "Die 10 Import-Dateien (2/2)",
        "Verknüpfung von Projekten & Ressourcen",
        [
            ("4 · Allocation & Tasks — das Herzstück: Zuweisungen + Plan/Ist-Werte;", 0),
            ("verbindet Projekt + Ressource + Task (← Ressourcen & Projekte)", 1),
            ("9 · Financials — Kosten je Projekt & Kategorie, Plan/Ist (← Projekte)", 0),
            ("7 · Sheets — Daten für bestehende Sheet-Templates; Zeilen werden angehängt", 0),
            ("Jede Datei hat ein festes Arbeitsblatt (Sheet) und einen klaren Zweck.", 0),
        ],
    ),
    (
        "Zusammenhänge der Dateien",
        "Drei Ebenen – was unten steht, braucht die Ebene darüber",
        [
            ("Ebene 1 · Fundament: Attributes (1) · Resources (2) · Skills (6)", 0),
            ("Ebene 2 · baut auf: Projects (3) · Advanced Rates (8) · Team (10) · Admin Times (5)", 0),
            ("Ebene 3 · verknüpft: Allocation (4) · Financials (9) · Sheets (7)", 0),
            ("Stolperstein: Eine Zuweisung (4) auf „Projekt X / Max M.“ scheitert,", 0),
            ("wenn Projekt X oder Max M. noch nicht existieren.", 1),
            ("Referenzierte Objekte müssen vorher angelegt sein.", 1),
        ],
    ),
    (
        "Die richtige Reihenfolge",
        "Die Nummerierung 1 → 10 bildet die Abhängigkeiten ab",
        [
            ("1. Attributes — Felder & Auswahllisten (keine Voraussetzung)", 0),
            ("2. Resources — Ressourcen-Pool (keine Voraussetzung)", 0),
            ("3. Projects — Container, nutzen Schritt 1 (braucht 1)", 0),
            ("4. Allocation & Tasks — Zuweisungen (braucht 2 & 3)", 0),
            ("5. Admin Times (braucht 2,1) · 6. Skills (unabhängig)", 0),
            ("7. Sheets (braucht 2,3) · 8. Advanced Rates (braucht 2)", 0),
            ("9. Financials (braucht 3,1) · 10. Team Resources (braucht 2)", 0),
            ("Faustregel: Stammdaten → Container → Verknüpfungen → Details", 0),
        ],
    ),
    (
        "Wie Tempus Zeilen wiedererkennt",
        "Der Schlüssel entscheidet: neu anlegen oder aktualisieren",
        [
            ("Attributes (1): Custom Field Name + Entity Type", 0),
            ("Resources (2): Resource Name (oder API External Id)", 0),
            ("Projects (3): Project Name (oder API External Id)", 0),
            ("Allocation (4): Project + Resource + Task + Plan Type + Allocation Type", 0),
            ("Praxis: Namen über alle Dateien identisch halten – sonst doppelte Datensätze.", 0),
            ("Projekt-Level ohne echten Task: eingebauten Task-Namen „Generic“ nutzen.", 0),
        ],
    ),
    (
        "Import Behavior",
        "Merge · Overwrite · Skip — pro Zeile steuerbar",
        [
            ("Merge — legt Neues an und aktualisiert Bestehendes;", 0),
            ("in Tempus gepflegte Werte bleiben erhalten. Standard für fast alles.", 1),
            ("Overwrite — ersetzt einen gleichnamigen Datensatz vollständig;", 0),
            ("bisherige Werte gehen verloren. Nur bewusst einsetzen.", 1),
            ("Skip — überspringt die komplette Zeile; nichts wird importiert.", 0),
            ("Im Zweifel immer Merge.", 0),
        ],
    ),
    (
        "Zeit-phasierte Spalten",
        "Datums-Spalten hinter den festen Spalten",
        [
            ("Resources (2): Monat — Stunden (Kapazität)", 0),
            ("Allocation (4): Projekt/Jahr/Quartal/Monat/Woche/Tag — Hours oder FTE", 0),
            ("Admin Times (5): Woche · Financials (9): Monat/Quartal/Jahr · Team (10): Woche/FTE", 0),
            ("Time Period muss zu den Spaltenköpfen passen (Week = immer Montag).", 0),
            ("Hours = Aufwand über Arbeitstage verteilt.", 0),
            ("FTE = Kapazitätsanteil über Arbeitstage repliziert.", 0),
        ],
    ),
    (
        "Schritt-für-Schritt & häufige Fehler",
        "Sauberer Erst-Import",
        [
            ("Ablauf: Attributes → Resources → Projects → Allocation → Details (5–10)", 0),
            ("Jede Datei separat hochladen und Rückmeldung prüfen.", 0),
            ("Häufige Fehler:", 0),
            ("Auswahlwert/Attribut fehlt → vorab in Attributes (1) anlegen", 1),
            ("Referenz fehlt → Reihenfolge einhalten (2 & 3 vor 4)", 1),
            ("Falsches Datumsformat / Spalten weichen ab → Vorlage unverändert nutzen", 1),
            ("Vorlage nicht erstes Tab → Datenblatt an erste Stelle ziehen", 1),
        ],
    ),
    (
        "Checkliste & nächste Schritte",
        "Vor jedem Import prüfen",
        [
            ("Richtige Vorlage · erstes Tab · Format .xlsx", 0),
            ("Original-Spaltenüberschriften unverändert · Pflichtfelder gefüllt", 0),
            ("Namen über alle Dateien identisch", 0),
            ("Reihenfolge 1 → 10 · Import Behavior bewusst (i. d. R. Merge)", 0),
            ("Datums-Spalten passen zu Time Period · Auswahlwerte existieren", 0),
            ("Erst Sandbox, dann Produktion", 0),
            ("Nächster Schritt: gemeinsamer Test-Import in der Sandbox (Schritt 1–4 live).", 0),
        ],
    ),
]


def main() -> int:
    prs = Presentation(str(MASTER))
    n_original = len(prs.slides._sldIdLst)

    # 1) Titelfolie
    s = clone_slide(prs, TPL_TITLE)
    texts = _shapes_with_text(s)
    # größte Schrift = Haupttitel
    title_shape = max(texts, key=lambda sh: (sh.text_frame.paragraphs[0].runs[0].font.size.pt
                                             if sh.text_frame.paragraphs[0].runs and sh.text_frame.paragraphs[0].runs[0].font.size else 0))
    set_title(title_shape, TITLE)

    # 2) Agenda
    s = clone_slide(prs, TPL_AGENDA)
    texts = _shapes_with_text(s)
    set_title(texts[0], "Agenda")
    if len(texts) > 1:
        set_body(texts[1], AGENDA)

    # 3) Inhaltsfolien
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for title, eyebrow, items in CONTENT:
        s = clone_slide(prs, TPL_CONTENT)
        shapes = list(s.shapes)
        title_sh = shapes[0]   # Titel oben
        body_sh = shapes[1]    # Body
        eyebrow_sh = shapes[3] if len(shapes) > 3 and shapes[3].has_text_frame else None
        set_title(title_sh, title)
        set_body(body_sh, items)
        if eyebrow_sh is not None:
            set_title(eyebrow_sh, eyebrow)
        # dekorative "Pins"-Navigationsgrafik (fremde Feature-Labels) entfernen
        for sh in shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                sh._element.getparent().remove(sh._element)

    # 4) Original-Folien entfernen
    delete_slides(prs, list(range(n_original)))

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    removed = cleanup_orphan_slide_parts(OUT)
    size_mb = OUT.stat().st_size / 1e6
    msg = f"✅ {len(prs.slides._sldIdLst)} Folien gespeichert: {OUT} ({size_mb:.1f} MB)"
    if removed:
        msg += f" · {removed} verwaiste Slide-Parts entfernt"
    print(msg)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
